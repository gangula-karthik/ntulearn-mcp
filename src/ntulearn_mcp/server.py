"""NTULearn MCP server — Blackboard Learn REST API wrapper."""

from __future__ import annotations

import asyncio
import base64
import contextlib
import json
import os
import re
import sys
from datetime import datetime, timedelta, timezone
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib.parse import unquote

from bs4 import BeautifulSoup
from dotenv import load_dotenv
from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import (
    GetPromptResult,
    ImageContent,
    Prompt,
    PromptArgument,
    PromptMessage,
    ReadResourceResult,
    Resource,
    ResourceTemplate,
    TextContent,
    TextResourceContents,
    Tool,
)

from ntulearn_mcp.cache import (
    delete_cached_cookie,
    read_cached_cookie,
    write_cached_cookie,
)
from ntulearn_mcp.client import BbRouterExpiredError, NTULearnClient
from ntulearn_mcp.cookie import read_bbrouter_cookie
from ntulearn_mcp.parsers import extract_all_files

load_dotenv()

# ---------------------------------------------------------------------------
# Config
# ---------------------------------------------------------------------------

BASE_URL = os.getenv("NTULEARN_BASE_URL", "https://ntulearn.ntu.edu.sg")
DOWNLOAD_DIR = Path(os.getenv("NTULEARN_DOWNLOAD_DIR", "./downloads"))

_TOOL_PREFIX = "ntulearn"

# Per-file and per-batch caps for read_file_content. download_bytes buffers the
# full response in memory, so without these a 200 MB attachment crashes the
# process and floods Claude's context with useless content.
_MAX_FILE_BYTES = 25 * 1024 * 1024
_MAX_TOTAL_BYTES = 40 * 1024 * 1024
# Per-sheet row cap for .xlsx extraction. Grade dumps and analytics exports
# can run to millions of rows; rendering all of them blows up the response.
_MAX_XLSX_ROWS_PER_SHEET = 1000

# Cap pages rendered as ImageContent in PDF vision mode. Each rendered page is
# ~2-3K vision tokens, so 50 pages ≈ 125K tokens — about the largest single
# attachment a user would realistically read in one tool call. Use the `pages`
# arg to step through bigger PDFs.
_MAX_PDF_PAGES_VISION = 50

# Cumulative PNG byte budget for PDF vision mode. claude.ai web caps a single
# tool result at ~1 MB; base64-inflated images push past that fast. Stop
# rendering once we've produced this many bytes of PNG. 800 KB leaves headroom
# for the structured payload + text and stays comfortably under the 1 MB cap
# even after base64 expansion (~33%).
_MAX_PDF_VISION_BYTES = 800 * 1024

# Render zoom for PDF vision. 1.3x ≈ 96 DPI is plenty for the API's 1568px
# downscale and roughly halves PNG byte size vs the previous 2.0x default.
_PDF_VISION_ZOOM = 1.3

# Pagination defaults / caps for list-returning tools.
_DEFAULT_LIMIT = 50
_MAX_LIMIT = 200
_MAX_DEPTH = 10

# Loose pattern for Blackboard course/content IDs (e.g. _12345_1, _abc-1_2).
_BB_ID_PATTERN = r"^[A-Za-z0-9_\-:]+$"

_TEXT_MIMETYPES = frozenset({
    "application/json",
    "application/xml",
    "application/javascript",
    "application/x-javascript",
    "application/x-yaml",
    "application/yaml",
    "application/ld+json",
    "application/x-sh",
})

_TEXT_EXTENSIONS = frozenset({
    "txt", "md", "markdown", "csv", "tsv", "json", "xml", "yaml", "yml",
    "html", "htm", "log", "py", "js", "ts", "rs", "go", "java", "c", "cpp",
    "h", "hpp", "sh", "bash", "zsh", "rb", "swift", "kt", "scala", "r",
    "ini", "toml", "cfg", "conf", "env",
})

_NO_COOKIE_MESSAGE = (
    "No NTULearn cookie found. Two options:\n"
    "  1. Log into https://ntulearn.ntu.edu.sg in a supported browser "
    "(Firefox works best on Windows; Chrome/Edge auto-read often needs "
    "admin), then restart your MCP host (e.g. Claude Desktop). This is "
    "the primary path — the server tries browsers first.\n"
    "  2. Set the NTULEARN_COOKIE env var manually as a fallback — see "
    "README for the DevTools cookie-copy steps. This always works but "
    "is only consulted when the browser auto-read fails."
)

# ---------------------------------------------------------------------------
# Server + client setup
# ---------------------------------------------------------------------------

app = Server("ntulearn_mcp")
_client: NTULearnClient | None = None


def _validate_cookie_value(value: str) -> str:
    """Reject cookie values that would let an attacker inject extra headers.

    A pathological value with CR/LF would smuggle Set-Cookie or other
    headers into the request. Realistic exposure is small (the user has
    to put the bad value in NTULEARN_COOKIE themselves) but cheap defence
    in depth.
    """
    if "\r" in value or "\n" in value or "\x00" in value:
        raise RuntimeError(
            "NTULEARN_COOKIE contains illegal control characters "
            "(CR/LF/NUL). Re-copy the cookie value from DevTools."
        )
    return value


def _resolve_cookie() -> str:
    """Resolve the BbRouter cookie.

    Resolution order (browser-first):
      1. ``read_bbrouter_cookie()`` — auto-read from a logged-in browser.
         This is the convenience path the MCP server exists for: zero
         config when the user is logged into NTULearn somewhere we can
         read. Bounded retry + backoff inside ``read_bbrouter_cookie``
         absorbs transient races. On success the value is mirrored to
         the OS keychain via ``write_cached_cookie`` so subsequent
         resolutions can fall back to it if the browser path later
         fails.
      2. ``NTULEARN_COOKIE`` env var — manual fallback. Used when the
         browser path can't deliver (Windows + Chrome/Edge ABE, no
         logged-in browser, headless environments). The env var is a
         safety net, not an override: a fresh browser read wins over a
         possibly-stale env value, which is what the user almost
         always wants when both are present.
      3. ``read_cached_cookie()`` — last-known-good value from the OS
         keychain. Catches the case where the browser fails right now
         AND the user hasn't seeded an env var, but a previous run did
         successfully read from the browser. Stretches a single
         successful read across the cookie's full lifetime.
      4. Raise — nothing left to try.

    Cache writes only happen on the browser path: an env var value is
    user-deliberate and not necessarily worth persisting (could be a
    one-off debug value), and a cache hit obviously doesn't need
    re-caching. The 401-retry path in ``_refresh_client`` invalidates
    the cache before re-resolving so we never loop on a dead value.
    """
    auto = read_bbrouter_cookie()
    if auto:
        # Persist for next time. Best-effort: a missing/broken keyring
        # backend just returns False and we move on.
        write_cached_cookie(auto)
        return _validate_cookie_value(auto)

    explicit = os.getenv("NTULEARN_COOKIE", "").strip()
    if explicit:
        return _validate_cookie_value(explicit)

    cached = read_cached_cookie()
    if cached:
        # Browser failed (transient race, ABE on Windows, browser
        # signed out, etc.) and there's no env var override. Cached
        # value is from the most recent successful read — likely still
        # valid since BbRouter cookies last days–weeks. If it's
        # expired, the 401-retry path in _refresh_client will nuke it
        # and force a re-resolve.
        return _validate_cookie_value(cached)

    raise RuntimeError(_NO_COOKIE_MESSAGE)


def get_client() -> NTULearnClient:
    global _client
    if _client is None:
        _client = NTULearnClient(BASE_URL, _resolve_cookie())
    return _client


async def _refresh_client() -> NTULearnClient:
    """Discard the current client, re-read the cookie, build a fresh client.

    Called after a 401 so an expired cookie can be transparently swapped
    for the fresh value the user's browser already has. Invalidates the
    keychain cache first — the cookie we just used produced a 401, so it's
    dead, and leaving it in the cache would let the next resolution loop
    back to the same dead value.
    """
    global _client
    if _client is not None:
        await _client.close()
        _client = None
    delete_cached_cookie()
    _client = NTULearnClient(BASE_URL, _resolve_cookie())
    return _client


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _strip_content(item: dict[str, Any]) -> dict[str, Any]:
    """Reduce a raw Blackboard content item to the fields we care about."""
    handler = item.get("contentHandler", {}) or {}
    description_raw = item.get("description", {}) or {}
    return {
        "id": item.get("id"),
        "title": item.get("title"),
        "contentHandlerId": handler.get("id"),
        "hasChildren": item.get("hasChildren", False),
        "description": description_raw.get("rawText") if isinstance(description_raw, dict) else description_raw,
        "modified": item.get("modified"),
    }


def _strip_calendar_item(item: dict[str, Any], course_id: str | None) -> dict[str, Any]:
    """Reduce a raw calendar item to fields useful to a student / agent.

    `dynamicCalendarItemProps` carries the interesting GradebookColumn metadata
    (eventType, gradable, attemptable) — flatten so callers don't have to dig.
    """
    dyn = item.get("dynamicCalendarItemProps") or {}
    return {
        "id": item.get("id"),
        "type": item.get("type"),
        "title": item.get("title"),
        "description": item.get("description"),
        "location": item.get("location"),
        "start": item.get("start"),
        "end": item.get("end"),
        "calendarName": item.get("calendarName"),
        "courseId": course_id,
        "eventType": dyn.get("eventType"),
        "gradable": dyn.get("gradable"),
        "attemptable": dyn.get("attemptable"),
    }


def _validate_iso8601(value: str, *, name: str) -> str:
    """Accept an ISO-8601 datetime string and return it.

    Blackboard's calendar API wants ``2026-05-09T00:00:00Z`` style. We round-trip
    through ``datetime.fromisoformat`` for a cheap sanity check — anything that
    parses survives. ``Z`` is normalised because ``fromisoformat`` only handles
    it natively from 3.11+, and we want a clear error rather than a 400 from
    Blackboard.
    """
    if not isinstance(value, str) or not value.strip():
        raise ValueError(f"{name} must be a non-empty ISO-8601 timestamp string")
    normalised = value.replace("Z", "+00:00") if value.endswith("Z") else value
    try:
        datetime.fromisoformat(normalised)
    except ValueError as e:
        raise ValueError(
            f"{name}={value!r} is not a valid ISO-8601 timestamp. "
            "Expected format like '2026-05-09T00:00:00Z'."
        ) from e
    return value


async def _resolve_content_files(
    client: NTULearnClient, course_id: str, content_id: str
) -> tuple[dict[str, Any], str | None, list[tuple[str, str | None]]]:
    """Return (item, handler_id, pairs) for a content item.

    `pairs` is a list of (url, filename) tuples for every file attached to the
    item. Filename may be None when the parser couldn't determine one.
    Empty pairs means no resolvable file (caller decides how to surface that).
    """
    item = await client.get_content_item(course_id, content_id)
    handler_id = (item.get("contentHandler") or {}).get("id")

    pairs: list[tuple[str, str | None]] = []

    if handler_id == "resource/x-bb-file":
        attachments = await client.get_attachments(course_id, content_id)
        for att in attachments:
            url = await client.get_attachment_download_url(course_id, content_id, att["id"])
            if url:
                pairs.append((url, att.get("fileName")))
    else:
        body = item.get("body") or ""
        files = extract_all_files(body)
        if not files:
            desc = item.get("description") or {}
            body2 = (desc.get("rawText") if isinstance(desc, dict) else desc) or ""
            files = extract_all_files(body2)
        for f in files:
            url = f.get("url")
            if url:
                pairs.append((url, f.get("filename")))

    return item, handler_id, pairs


def _file_extension(filename: str) -> str:
    if "." not in filename:
        return ""
    return filename.rpartition(".")[2].lower()


def _strip_html(value: Any) -> str:
    """Strip HTML tags and collapse whitespace from a possibly-HTML string.

    Used by both the file-content text path (HTML pages, e.g. course handouts saved
    as .html) and the announcements path (Blackboard stores rich-text bodies as
    HTML in `body.rawText`). Centralised so both call sites stay consistent.

    Accepts non-strings for caller convenience: ``None`` and other falsy values
    return ``""``; non-string truthy values are coerced via ``str()``.
    """
    if not value:
        return ""
    if not isinstance(value, str):
        value = str(value)
    text = BeautifulSoup(value, "html.parser").get_text(separator="\n")
    return "\n".join(
        line for line in (segment.strip() for segment in text.splitlines()) if line
    )


def _parse_content_type(content_type: str | None) -> tuple[str, str | None]:
    """Return (mime, charset) from a Content-Type header value."""
    if not content_type:
        return "", None
    parts = [p.strip() for p in content_type.split(";")]
    mime = parts[0].lower()
    charset = None
    for p in parts[1:]:
        if p.lower().startswith("charset="):
            charset = p.split("=", 1)[1].strip().strip('"').strip("'")
    return mime, charset


def _classify_kind(filename: str, content_type: str | None) -> str:
    """Return 'pdf', 'docx', 'pptx', 'xlsx', 'text', or 'binary'.

    Filename extension wins over content-type. Blackboard's bbcswebdav often
    serves files as application/octet-stream regardless of actual format, so
    trusting the header alone misclassifies most attachments.
    """
    ext = _file_extension(filename)
    if ext == "pdf":
        return "pdf"
    if ext == "docx":
        return "docx"
    if ext == "pptx":
        return "pptx"
    if ext == "xlsx":
        return "xlsx"
    if ext in _TEXT_EXTENSIONS:
        return "text"

    mime, _ = _parse_content_type(content_type)
    if mime == "application/pdf":
        return "pdf"
    # OOXML MIMEs are unwieldy
    # (e.g. application/vnd.openxmlformats-officedocument.wordprocessingml.document);
    # substring match is robust against minor spelling variants from servers.
    if "wordprocessingml" in mime:
        return "docx"
    if "presentationml" in mime:
        return "pptx"
    if "spreadsheetml" in mime:
        return "xlsx"
    if mime.startswith("text/"):
        return "text"
    if mime in _TEXT_MIMETYPES:
        return "text"
    return "binary"


def _format_bytes(n: int) -> str:
    if n < 1024:
        return f"{n} B"
    if n < 1024 * 1024:
        return f"{n / 1024:.1f} KB"
    return f"{n / (1024 * 1024):.1f} MB"


def _extract_content(
    filename: str, content_type: str | None, content_bytes: bytes
) -> dict[str, Any]:
    """Detect file kind and extract text. Returns one entry for the response payload.

    Pure function — no I/O. Caller handles download retries and size caps.
    """
    size = len(content_bytes)
    kind = _classify_kind(filename, content_type)

    if kind == "pdf":
        from pypdf import PdfReader
        from pypdf.errors import PyPdfError

        try:
            reader = PdfReader(BytesIO(content_bytes))
            if reader.is_encrypted:
                # Many "encrypted" Blackboard PDFs unlock with an empty password.
                try:
                    reader.decrypt("")
                except Exception:
                    return {
                        "filename": filename,
                        "kind": "pdf",
                        "error": "PDF is password-protected. Cannot extract text.",
                        "sizeBytes": size,
                        "contentType": content_type,
                    }
            page_count = len(reader.pages)
            text = "\n\n".join((page.extract_text() or "") for page in reader.pages)
        except (PyPdfError, Exception) as e:
            return {
                "filename": filename,
                "kind": "pdf",
                "error": f"PDF extraction failed: {e}",
                "sizeBytes": size,
                "contentType": content_type,
            }

        out: dict[str, Any] = {
            "filename": filename,
            "kind": "pdf",
            "text": text,
            "pageCount": page_count,
            "sizeBytes": size,
            "contentType": content_type,
        }
        if not text.strip():
            out["warning"] = (
                "PDF appears to contain no extractable text "
                "(likely scanned images)."
            )
        return out

    if kind == "docx":
        return _extract_docx(filename, content_type, content_bytes, size)

    if kind == "pptx":
        return _extract_pptx(filename, content_type, content_bytes, size)

    if kind == "xlsx":
        return _extract_xlsx(filename, content_type, content_bytes, size)

    if kind == "text":
        _, charset = _parse_content_type(content_type)
        text: str | None = None
        for enc in (charset, "utf-8", "latin-1"):
            if not enc:
                continue
            try:
                text = content_bytes.decode(enc)
                break
            except (UnicodeDecodeError, LookupError):
                continue
        if text is None:
            text = content_bytes.decode("utf-8", errors="replace")

        ext = _file_extension(filename)
        mime, _ = _parse_content_type(content_type)
        if ext in {"html", "htm"} or mime == "text/html":
            text = _strip_html(text)

        return {
            "filename": filename,
            "kind": "text",
            "text": text,
            "sizeBytes": size,
            "contentType": content_type,
        }

    return {
        "filename": filename,
        "kind": "binary",
        "error": (
            f"Binary file ({content_type or 'unknown type'}, {_format_bytes(size)}). "
            "Cannot extract text. Use ntulearn_download_file to save it locally."
        ),
        "sizeBytes": size,
        "contentType": content_type,
    }


@contextlib.contextmanager
def _suppress_stdout_fd():
    """Redirect raw stdout fd (1) to /dev/null for the duration of the block.

    MuPDF emits warnings/errors from C, bypassing any Python-level
    ``sys.stdout`` redirection. Anything that reaches fd 1 corrupts the
    MCP JSON-RPC frame the host is parsing. Belt-and-suspenders alongside
    ``fitz.TOOLS.mupdf_display_errors(False)``.
    """
    sys.stdout.flush()
    devnull_fd = os.open(os.devnull, os.O_WRONLY)
    saved_fd = os.dup(1)
    try:
        os.dup2(devnull_fd, 1)
        yield
    finally:
        sys.stdout.flush()
        os.dup2(saved_fd, 1)
        os.close(saved_fd)
        os.close(devnull_fd)


def _extract_pdf_vision(
    filename: str,
    content_type: str | None,
    content_bytes: bytes,
    size: int,
    pages_filter: set[int] | None,
) -> dict[str, Any]:
    """Render PDF pages to PNG + extract text per page via PyMuPDF.

    Mirrors what Claude.ai does when a user attaches a PDF: per-page text
    AND a rendered image, so the model can see diagrams / equations /
    scanned content that pure-text extractors miss. Each rendered page
    costs roughly 1-2K vision tokens; rendering stops at whichever of
    _MAX_PDF_PAGES_VISION or _MAX_PDF_VISION_BYTES is hit first.

    Returns a dict with the structured-content fields and a private
    ``_images`` list of ``(label, png_bytes)`` tuples for the caller to
    emit as ``ImageContent`` blocks. The structured payload itself does
    not include image bytes.
    """
    try:
        import fitz  # pymupdf
    except ImportError as e:  # pragma: no cover — dep declared in pyproject
        return {
            "filename": filename,
            "kind": "pdf",
            "error": f"pymupdf not available: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }

    # Silence MuPDF's Python-visible diagnostic stream. fd-level redirect
    # catches anything the C layer emits regardless. Both are necessary —
    # invalid-PDF errors fire during open, well before the render loop.
    try:
        fitz.TOOLS.mupdf_display_errors(False)
    except Exception:
        pass

    try:
        with _suppress_stdout_fd():
            doc = fitz.open(stream=content_bytes, filetype="pdf")
    except Exception as e:
        return {
            "filename": filename,
            "kind": "pdf",
            "error": f"PDF could not be opened: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }

    try:
        if doc.is_encrypted and not doc.authenticate(""):
            return {
                "filename": filename,
                "kind": "pdf",
                "error": "PDF is password-protected. Cannot extract text.",
                "sizeBytes": size,
                "contentType": content_type,
            }

        total_pages = doc.page_count
        if pages_filter is None:
            page_indices = list(range(total_pages))
        else:
            # 1-indexed in the API, 0-indexed for PyMuPDF.
            page_indices = sorted(p - 1 for p in pages_filter if 1 <= p <= total_pages)

        truncation_reason: str | None = None
        if len(page_indices) > _MAX_PDF_PAGES_VISION:
            page_indices = page_indices[:_MAX_PDF_PAGES_VISION]
            truncation_reason = "page_cap"

        text_parts: list[str] = []
        images: list[tuple[str, bytes]] = []
        rendered_indices: list[int] = []
        cumulative_bytes = 0
        zoom = fitz.Matrix(_PDF_VISION_ZOOM, _PDF_VISION_ZOOM)
        with _suppress_stdout_fd():
            for idx in page_indices:
                page = doc[idx]
                page_text = page.get_text()
                pix = page.get_pixmap(matrix=zoom, alpha=False)
                png_bytes = pix.tobytes("png")
                if (
                    cumulative_bytes + len(png_bytes) > _MAX_PDF_VISION_BYTES
                    and rendered_indices
                ):
                    truncation_reason = "byte_budget"
                    break
                cumulative_bytes += len(png_bytes)
                text_parts.append(f"## Page {idx + 1}\n\n{page_text}")
                images.append((f"{filename} · page {idx + 1}", png_bytes))
                rendered_indices.append(idx)

        text = "\n\n".join(text_parts)
        truncated_pages = len(page_indices) - len(rendered_indices)
        out: dict[str, Any] = {
            "filename": filename,
            "kind": "pdf",
            "text": text,
            "pageCount": total_pages,
            "pagesRendered": [i + 1 for i in rendered_indices],
            "_images": images,
            "sizeBytes": size,
            "contentType": content_type,
        }
        if truncation_reason is not None:
            out["truncatedPages"] = truncated_pages
            out["truncationReason"] = truncation_reason
            if truncation_reason == "byte_budget":
                next_page = rendered_indices[-1] + 2 if rendered_indices else 1
                out["warning"] = (
                    f"Rendered {len(rendered_indices)} page(s) before hitting the "
                    f"~{_MAX_PDF_VISION_BYTES // 1024} KB image-bytes budget. "
                    f"{truncated_pages} page(s) of the requested range were not rendered. "
                    f"Call again with pages='{next_page}-' to continue, or pass "
                    f"mode='text' for a cheaper text-only pass."
                )
            else:
                out["warning"] = (
                    f"Requested range exceeds the {_MAX_PDF_PAGES_VISION}-page render cap; "
                    f"only the first {_MAX_PDF_PAGES_VISION} page(s) were rendered. "
                    "Use the `pages` arg to step through the rest."
                )
        return out
    except Exception as e:
        return {
            "filename": filename,
            "kind": "pdf",
            "error": f"PDF vision extraction failed: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }
    finally:
        doc.close()


def _extract_docx(
    filename: str, content_type: str | None, content_bytes: bytes, size: int
) -> dict[str, Any]:
    """Extract paragraph + table text from a .docx (OOXML) file."""
    from io import BytesIO

    try:
        from docx import Document
        from docx.opc.exceptions import PackageNotFoundError
    except ImportError as e:  # pragma: no cover — dep declared in pyproject
        return {
            "filename": filename,
            "kind": "docx",
            "error": f"python-docx not available: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }

    try:
        doc = Document(BytesIO(content_bytes))
    except PackageNotFoundError as e:
        return {
            "filename": filename,
            "kind": "docx",
            "error": f"Not a valid .docx file: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }
    except Exception as e:
        return {
            "filename": filename,
            "kind": "docx",
            "error": f"DOCX extraction failed: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }

    paragraphs = [p.text for p in doc.paragraphs if p.text and p.text.strip()]
    tables_text: list[str] = []
    for table in doc.tables:
        rows: list[str] = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                rows.append("\t".join(cells))
        if rows:
            tables_text.append("\n".join(rows))

    parts = ["\n\n".join(paragraphs)] if paragraphs else []
    if tables_text:
        parts.append("## Tables\n\n" + "\n\n".join(tables_text))
    text = "\n\n".join(parts)

    return {
        "filename": filename,
        "kind": "docx",
        "text": text,
        "paragraphCount": len(paragraphs),
        "tableCount": len(doc.tables),
        "sizeBytes": size,
        "contentType": content_type,
    }


def _extract_pptx(
    filename: str, content_type: str | None, content_bytes: bytes, size: int
) -> dict[str, Any]:
    """Extract slide text + speaker notes from a .pptx (OOXML) file."""
    from io import BytesIO

    try:
        from pptx import Presentation
    except ImportError as e:  # pragma: no cover — dep declared in pyproject
        return {
            "filename": filename,
            "kind": "pptx",
            "error": f"python-pptx not available: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }

    try:
        prs = Presentation(BytesIO(content_bytes))
    except Exception as e:
        return {
            "filename": filename,
            "kind": "pptx",
            "error": f"PPTX extraction failed: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }

    slide_blocks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        shape_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = shape.text_frame.text
                if t and t.strip():
                    shape_texts.append(t)
            elif shape.has_table:
                rows: list[str] = []
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        rows.append("\t".join(cells))
                if rows:
                    shape_texts.append("\n".join(rows))

        block = f"## Slide {idx}"
        if shape_texts:
            block += "\n\n" + "\n\n".join(shape_texts)

        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            notes = notes_tf.text if notes_tf else ""
            if notes and notes.strip():
                block += f"\n\nSpeaker notes:\n{notes}"

        slide_blocks.append(block)

    text = "\n\n".join(slide_blocks)
    return {
        "filename": filename,
        "kind": "pptx",
        "text": text,
        "slideCount": len(prs.slides),
        "sizeBytes": size,
        "contentType": content_type,
    }


def _extract_xlsx(
    filename: str, content_type: str | None, content_bytes: bytes, size: int
) -> dict[str, Any]:
    """Extract row data from all sheets of an .xlsx (OOXML) workbook."""
    from io import BytesIO

    try:
        from openpyxl import load_workbook
    except ImportError as e:  # pragma: no cover — dep declared in pyproject
        return {
            "filename": filename,
            "kind": "xlsx",
            "error": f"openpyxl not available: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }

    try:
        wb = load_workbook(BytesIO(content_bytes), data_only=True, read_only=True)
    except Exception as e:
        return {
            "filename": filename,
            "kind": "xlsx",
            "error": f"XLSX extraction failed: {e}",
            "sizeBytes": size,
            "contentType": content_type,
        }

    truncated = False
    sheet_blocks: list[str] = []
    try:
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_text: list[str] = []
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                if all(c is None or (isinstance(c, str) and not c.strip()) for c in row):
                    continue
                if row_count >= _MAX_XLSX_ROWS_PER_SHEET:
                    truncated = True
                    rows_text.append(
                        f"... (additional rows truncated; cap is "
                        f"{_MAX_XLSX_ROWS_PER_SHEET} rows per sheet)"
                    )
                    break
                cells = ["" if c is None else str(c) for c in row]
                rows_text.append("\t".join(cells))
                row_count += 1
            block = f"## Sheet: {sheet_name}"
            if rows_text:
                block += "\n\n" + "\n".join(rows_text)
            sheet_blocks.append(block)
    finally:
        wb.close()

    text = "\n\n".join(sheet_blocks)
    out: dict[str, Any] = {
        "filename": filename,
        "kind": "xlsx",
        "text": text,
        "sheetCount": len(sheet_blocks),
        "sizeBytes": size,
        "contentType": content_type,
    }
    if truncated:
        out["warning"] = (
            f"One or more sheets exceeded {_MAX_XLSX_ROWS_PER_SHEET} rows "
            "and were truncated. Use download_file for the full data."
        )
    return out


# ---------------------------------------------------------------------------
# Pagination + response-format helpers
# ---------------------------------------------------------------------------

def _slice_with_pagination(
    items: list[Any], offset: int, limit: int
) -> tuple[list[Any], dict[str, Any]]:
    """Slice items[offset:offset+limit] and return pagination metadata.

    The underlying Blackboard client paginates internally and returns full
    result sets, so this is caller-side slicing — its job is to keep the
    LLM's context bounded, not to reduce upstream load.
    """
    total = len(items)
    end = min(total, offset + limit)
    page = items[offset:end]
    next_offset = end if end < total else None
    return page, {
        "total": total,
        "count": len(page),
        "offset": offset,
        "limit": limit,
        "hasMore": next_offset is not None,
        "nextOffset": next_offset,
    }


def _resolve_pagination_args(args: dict[str, Any]) -> tuple[int, int]:
    """Validate and clamp limit/offset args. Defaults: offset=0, limit=_DEFAULT_LIMIT."""
    offset = int(args.get("offset", 0))
    limit = int(args.get("limit", _DEFAULT_LIMIT))
    if offset < 0:
        raise ValueError("offset must be >= 0")
    if limit < 1:
        raise ValueError("limit must be >= 1")
    if limit > _MAX_LIMIT:
        raise ValueError(f"limit must be <= {_MAX_LIMIT}")
    return offset, limit


def _resolve_response_format(args: dict[str, Any]) -> str:
    """Return 'json' or 'markdown'. Default: 'json' (machine-readable)."""
    fmt = str(args.get("response_format", "json")).lower()
    if fmt not in ("json", "markdown"):
        raise ValueError("response_format must be 'json' or 'markdown'")
    return fmt


def _resolve_pdf_mode(args: dict[str, Any]) -> str:
    """Return 'text' or 'vision'. Default: 'text'.

    Most PDF reads on NTULearn are text-extraction questions ("what's the due
    date in this brief?", "summarise this reading"); rendering every page as an
    image bloats the response past MCP's 1 MB cap for typical lecture decks.
    Vision is opt-in for diagram/equation-heavy pages — pair with `pages` to
    restrict to the specific page(s) you need.

    'auto' is accepted as an alias for 'text' for backwards compatibility with
    earlier callers; treating auto as text is the safe default.
    """
    mode = str(args.get("mode", "text")).lower()
    if mode == "auto":
        mode = "text"
    if mode not in ("text", "vision"):
        raise ValueError("mode must be 'text' or 'vision'")
    return mode


def _parse_page_range(spec: Any) -> set[int] | None:
    """Parse a page-range spec into a set of 1-indexed page numbers, or None for all.

    Accepts strings like "1-10", "1,3,5", "1-5,8,10-12". Whitespace is
    ignored. Returns None when spec is missing or empty so the caller can
    treat "no filter" as "render every page".
    """
    if spec is None:
        return None
    s = str(spec).strip()
    if not s:
        return None
    pages: set[int] = set()
    for token in s.split(","):
        t = token.strip()
        if not t:
            continue
        if "-" in t:
            lo_s, _, hi_s = t.partition("-")
            try:
                lo, hi = int(lo_s.strip()), int(hi_s.strip())
            except ValueError:
                raise ValueError(f"Invalid page range token: {t!r}")
            if lo < 1 or hi < lo:
                raise ValueError(f"Invalid page range: {t!r} (need 1 <= lo <= hi)")
            pages.update(range(lo, hi + 1))
        else:
            try:
                n = int(t)
            except ValueError:
                raise ValueError(f"Invalid page number: {t!r}")
            if n < 1:
                raise ValueError(f"Page numbers are 1-indexed; got {n}")
            pages.add(n)
    return pages or None


def _emit(
    payload: dict[str, Any], text: str | None = None
) -> tuple[list[TextContent], dict[str, Any]]:
    """Return the (unstructured, structured) tuple MCP expects when both are present.

    `payload` is the JSON-serialisable structured content (validated against
    outputSchema by MCP). `text` is the rendered display text — defaults to a
    pretty-printed JSON copy of the payload for json mode.
    """
    if text is None:
        text = json.dumps(payload, indent=2)
    return [TextContent(type="text", text=text)], payload


def _md_pagination_footer(meta: dict[str, Any]) -> str:
    """Trailing line summarising pagination for markdown views."""
    if meta["hasMore"]:
        return (
            f"\n_Showing {meta['count']} of {meta['total']} "
            f"(offset {meta['offset']}). Pass offset={meta['nextOffset']} for the next page._"
        )
    return f"\n_Showing all {meta['total']}._"


# ---------------------------------------------------------------------------
# Tool definitions
# ---------------------------------------------------------------------------

# Common JSON-schema fragments for reuse.
_COURSE_ID_SCHEMA = {
    "type": "string",
    "description": "Blackboard course ID (e.g. _12345_1)",
    "minLength": 1,
    "maxLength": 200,
    "pattern": _BB_ID_PATTERN,
}

_CONTENT_ID_SCHEMA = {
    "type": "string",
    "description": "Content item ID (e.g. _67890_1)",
    "minLength": 1,
    "maxLength": 200,
    "pattern": _BB_ID_PATTERN,
}

_LIMIT_SCHEMA = {
    "type": "integer",
    "description": (
        f"Max items to return per call (1–{_MAX_LIMIT}, default {_DEFAULT_LIMIT}). "
        "Use small values to keep results within the LLM's context."
    ),
    "minimum": 1,
    "maximum": _MAX_LIMIT,
    "default": _DEFAULT_LIMIT,
}

_OFFSET_SCHEMA = {
    "type": "integer",
    "description": (
        "Number of items to skip for pagination. "
        "Use the nextOffset value returned by a previous call to walk pages."
    ),
    "minimum": 0,
    "default": 0,
}

_RESPONSE_FORMAT_SCHEMA = {
    "type": "string",
    "description": (
        "'json' returns a structured payload (default, recommended for agents). "
        "'markdown' returns a human-readable summary."
    ),
    "enum": ["json", "markdown"],
    "default": "json",
}

_PAGINATION_OUTPUT_FIELDS = {
    "total": {"type": "integer"},
    "count": {"type": "integer"},
    "offset": {"type": "integer"},
    "limit": {"type": "integer"},
    "hasMore": {"type": "boolean"},
    "nextOffset": {"type": ["integer", "null"]},
}

_CONTENT_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "id": {"type": ["string", "null"]},
        "title": {"type": ["string", "null"]},
        "contentHandlerId": {"type": ["string", "null"]},
        "hasChildren": {"type": "boolean"},
        "description": {"type": ["string", "null"]},
        "modified": {"type": ["string", "null"]},
    },
}

_ANNOUNCEMENT_SCHEMA = {
    "type": "object",
    "properties": {
        "id": {"type": ["string", "null"]},
        "courseId": {"type": ["string", "null"]},
        "title": {"type": ["string", "null"]},
        "body": {"type": ["string", "null"]},
        "created": {"type": ["string", "null"]},
        "modified": {"type": ["string", "null"]},
        "available": {"type": ["string", "null"]},
    },
}

_CALENDAR_ITEM_TYPES = ("Course", "GradebookColumn", "Institution", "OfficeHours", "Personal")

_CALENDAR_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "id": {"type": ["string", "null"]},
        "type": {"type": ["string", "null"]},
        "title": {"type": ["string", "null"]},
        "description": {"type": ["string", "null"]},
        "location": {"type": ["string", "null"]},
        "start": {"type": ["string", "null"]},
        "end": {"type": ["string", "null"]},
        "calendarName": {"type": ["string", "null"]},
        "courseId": {"type": ["string", "null"]},
        "eventType": {"type": ["string", "null"]},
        "gradable": {"type": ["boolean", "null"]},
        "attemptable": {"type": ["boolean", "null"]},
    },
}

_GRADEBOOK_COLUMN_SCHEMA = {
    "type": "object",
    "properties": {
        "id": {"type": ["string", "null"]},
        "courseId": {"type": ["string", "null"]},
        "name": {"type": ["string", "null"]},
        "displayName": {"type": ["string", "null"]},
        "possible": {"type": ["number", "null"]},
        "available": {"type": ["string", "null"]},
        "contentId": {"type": ["string", "null"]},
        "score": {"type": ["number", "string", "null"]},
        "grade": {"type": ["string", "null"]},
        "status": {"type": ["string", "null"]},
    },
}

_FILE_INFO_SCHEMA = {
    "type": "object",
    "properties": {
        "url": {"type": ["string", "null"]},
        "filename": {"type": ["string", "null"]},
        "mimeType": {"type": ["string", "null"]},
        "link_text": {"type": ["string", "null"]},
        # download_file results
        "localPath": {"type": "string"},
        "sizeBytes": {"type": "integer"},
        # read_file_content results
        "kind": {"type": "string"},
        "text": {"type": "string"},
        "pageCount": {"type": "integer"},
        "pagesRendered": {"type": "array", "items": {"type": "integer"}},
        "truncatedPages": {"type": "integer"},
        "truncationReason": {"type": "string", "enum": ["byte_budget", "page_cap"]},
        "paragraphCount": {"type": "integer"},
        "tableCount": {"type": "integer"},
        "slideCount": {"type": "integer"},
        "sheetCount": {"type": "integer"},
        "contentType": {"type": ["string", "null"]},
        "warning": {"type": "string"},
        "error": {"type": "string"},
        "reason": {"type": "string"},
    },
}


# ---------------------------------------------------------------------------
# New-tool schema fragments (v0.3) — shared with the tool handlers in
# ntulearn_mcp.handlers (WT-B). server.py only needs the schemas here to
# expose the tool surface; the handler bodies are imported lazily at call time
# so this module imports cleanly before handlers.py lands.
# ---------------------------------------------------------------------------

_COURSE_IDS_SCHEMA = {
    "type": "array",
    "description": (
        "Optional list of course IDs to scope to. Omit to fan out across all "
        "available enrolled courses."
    ),
    "items": {"type": "string", "pattern": _BB_ID_PATTERN},
    "maxItems": _MAX_LIMIT,
}

_MESSAGE_FOLDER_SCHEMA = {
    "type": "string",
    "description": (
        "Mailbox folder to read. Blackboard exposes 'inbox' (received) and "
        "'sent' (outbox)."
    ),
    "enum": ["inbox", "sent"],
    "default": "inbox",
}

_SINCE_SCHEMA = {
    "type": "string",
    "description": (
        "Optional ISO-8601 cutoff (e.g. '2026-05-09T00:00:00Z'). Only items "
        "created on/after this time are included."
    ),
}

_UNTIL_SCHEMA = {
    "type": "string",
    "description": (
        "Optional ISO-8601 end of the window (e.g. '2026-05-23T00:00:00Z')."
    ),
}

_QUERY_SCHEMA = {
    "type": "string",
    "description": "Search term (case-insensitive substring)",
    "minLength": 1,
    "maxLength": 200,
}

_MAX_DEPTH_SCHEMA = {
    "type": "integer",
    "description": f"Maximum recursion depth (capped at {_MAX_DEPTH}).",
    "default": 5,
    "minimum": 1,
    "maximum": _MAX_DEPTH,
}

_MESSAGE_SCHEMA = {
    "type": "object",
    "properties": {
        "id": {"type": ["string", "null"]},
        "subject": {"type": ["string", "null"]},
        "body": {"type": ["string", "null"]},
        "senderName": {"type": ["string", "null"]},
        "senderId": {"type": ["string", "null"]},
        "createdAt": {"type": ["string", "null"]},
        "read": {"type": "boolean"},
        "recipients": {"type": "array", "items": {"type": "object"}},
    },
}

_USER_SCHEMA = {
    "type": "object",
    "properties": {
        "id": {"type": ["string", "null"]},
        "name": {"type": ["string", "null"]},
        "role": {"type": ["string", "null"]},
        "userName": {"type": ["string", "null"]},
    },
}

_GROUP_SCHEMA = {
    "type": "object",
    "properties": {
        "id": {"type": ["string", "null"]},
        "name": {"type": ["string", "null"]},
        "description": {"type": ["string", "null"]},
        "available": {"type": ["string", "null"]},
    },
}

_ATTEMPT_SCHEMA = {
    "type": "object",
    "properties": {
        "id": {"type": ["string", "null"]},
        "userId": {"type": ["string", "null"]},
        "userName": {"type": ["string", "null"]},
        "status": {"type": ["string", "null"]},
        "score": {"type": ["number", "string", "null"]},
        "feedback": {"type": ["string", "null"]},
        "cumulatedScore": {"type": ["number", "string", "null"]},
        "createdAt": {"type": ["string", "null"]},
    },
}


@app.list_tools()
async def list_tools() -> list[Tool]:
    return [
        Tool(
            name=f"{_TOOL_PREFIX}_list_courses",
            description=(
                "List courses the current user is enrolled in on NTULearn. "
                "By default returns only active/available courses. "
                "Set include_disabled=true to also include unavailable ones. "
                "Supports pagination via limit/offset."
            ),
            annotations={
                "title": "List my NTULearn courses",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "include_disabled": {
                        "type": "boolean",
                        "description": "Include courses where availability.available != 'Yes'",
                        "default": False,
                    },
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "courses": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "courseId": {"type": "string"},
                                "title": {"type": "string"},
                                "available": {"type": "string"},
                                "lastAccessed": {"type": ["string", "null"]},
                            },
                            "required": ["courseId", "title"],
                        },
                    },
                    **_PAGINATION_OUTPUT_FIELDS,
                },
                "required": ["courses", "total", "count", "offset", "limit", "hasMore"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_get_course_contents",
            description=(
                "Walk a course's content tree. Omit parent_id to get the top-level "
                "items (folders, documents, links, assignments); pass parent_id of a "
                "folder/lesson where hasChildren=true to drill into its children. "
                "Supports pagination."
            ),
            annotations={
                "title": "Get course contents (root or folder children)",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "parent_id": {
                        **_CONTENT_ID_SCHEMA,
                        "description": (
                            "Optional content item ID of a folder/lesson to list "
                            "children of. Omit to list the course's top-level items."
                        ),
                    },
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "items": {"type": "array", "items": _CONTENT_ITEM_SCHEMA},
                    **_PAGINATION_OUTPUT_FIELDS,
                },
                "required": ["items", "total", "count", "offset", "limit", "hasMore"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_search_course_content",
            description=(
                "Recursively search a course's entire content tree for items matching a query. "
                "Matches on title or description (case-insensitive substring). "
                "Returns matched items with their full breadcrumb path."
            ),
            annotations={
                "title": "Search course content",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "query": {
                        "type": "string",
                        "description": "Search term (case-insensitive substring)",
                        "minLength": 1,
                        "maxLength": 200,
                    },
                    "max_depth": {
                        "type": "integer",
                        "description": f"Maximum recursion depth (default 5, capped at {_MAX_DEPTH})",
                        "default": 5,
                        "minimum": 1,
                        "maximum": _MAX_DEPTH,
                    },
                    "max_results": {
                        "type": "integer",
                        "description": f"Maximum number of matching items to return (default 50, capped at {_MAX_LIMIT})",
                        "default": 50,
                        "minimum": 1,
                        "maximum": _MAX_LIMIT,
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id", "query"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "matches": {
                        "type": "array",
                        "items": {
                            **_CONTENT_ITEM_SCHEMA,
                            "properties": {
                                **_CONTENT_ITEM_SCHEMA["properties"],
                                "breadcrumb": {"type": "array", "items": {"type": "string"}},
                            },
                        },
                    },
                    "count": {"type": "integer"},
                },
                "required": ["matches", "count"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_download_file",
            description=(
                "Download every file attached to a Blackboard content item to local "
                "disk. Handles both resource/x-bb-file (attachment API) and "
                "resource/x-bb-document (HTML body with bbcswebdav links) handler types. "
                "Pass destination_dir to target a specific folder — useful for "
                "organising a semester (e.g. destination_dir='~/NTU/y3s1/sc2002/week 8/'). "
                "Returns saved files with their resolved local paths and sizes. "
                "Use ntulearn_read_file_content if you want to inspect the content "
                "inline rather than saving to disk."
            ),
            annotations={
                "title": "Download files to local disk",
                "readOnlyHint": False,
                "destructiveHint": False,
                "idempotentHint": False,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "content_id": _CONTENT_ID_SCHEMA,
                    "destination_dir": {
                        "type": "string",
                        "description": (
                            "Optional target directory. Accepts absolute paths and "
                            "`~`-prefixed paths (e.g. '~/NTU/y3s1/sc2002/week 8/'). "
                            "Created if missing. Defaults to NTULEARN_DOWNLOAD_DIR "
                            "env var, or ./downloads/ if unset."
                        ),
                        "minLength": 1,
                        "maxLength": 1024,
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id", "content_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "contentId": {"type": ["string", "null"]},
                    "title": {"type": ["string", "null"]},
                    "contentHandlerId": {"type": ["string", "null"]},
                    "files": {"type": "array", "items": _FILE_INFO_SCHEMA},
                    "destinationDir": {"type": "string"},
                    "error": {"type": "string"},
                },
                "required": ["files"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_read_file_content",
            description=(
                "Read the content of files attached to a Blackboard content item, "
                "returned inline (no local-filesystem hop). Use this to ask questions "
                "about lecture material — ntulearn_download_file is for users who "
                "actually want the bytes on disk. "
                "PDFs default to text mode (via pypdf — cheap and almost always "
                "what you want for written content). Pass mode='vision' for "
                "diagram-, equation-, or screenshot-heavy pages; pair with "
                "pages='5' or pages='1-3' to keep the payload under MCP's 1 MB "
                "cap (~3K vision tokens per page). For multi-page diagram-heavy "
                "decks, prefer ntulearn_download_file plus drag-and-drop into "
                "claude.ai. "
                "Also supports Microsoft Office formats (.docx, .pptx with speaker "
                "notes, .xlsx with all sheets) and text-like files (txt, md, csv, "
                "json, xml, html with tags stripped, code files). "
                "Other binaries (images, video, audio, archives, legacy .doc/.ppt/.xls) "
                "are listed under `skipped` — fall back to ntulearn_download_file for those. "
                f"Per-file cap 25 MB, batch cap 40 MB, vision cap "
                f"{_MAX_PDF_PAGES_VISION} rendered pages; oversized files are skipped."
            ),
            annotations={
                "title": "Read file content inline",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "content_id": _CONTENT_ID_SCHEMA,
                    "mode": {
                        "type": "string",
                        "description": (
                            "PDF handling. 'text' (default) extracts text via "
                            "pypdf — cheap, fits MCP's payload budget. 'vision' "
                            "additionally renders each page as an image with "
                            "PyMuPDF (~3K vision tokens per page); use for "
                            "diagram/equation/handwritten content, ideally with "
                            "a narrow `pages` range to stay under the 1 MB cap. "
                            "Ignored for non-PDF files."
                        ),
                        "enum": ["text", "vision", "auto"],
                        "default": "text",
                    },
                    "pages": {
                        "type": "string",
                        "description": (
                            "Optional page range for PDFs (1-indexed, inclusive). "
                            "Examples: '1-10', '3', '1,3,5', '1-5,8,10-12'. Omit "
                            f"to read all pages (vision mode capped at "
                            f"{_MAX_PDF_PAGES_VISION} rendered pages). Especially "
                            "useful with mode='vision' to keep the response "
                            "under MCP's 1 MB cap."
                        ),
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id", "content_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "contentId": {"type": ["string", "null"]},
                    "title": {"type": ["string", "null"]},
                    "contentHandlerId": {"type": ["string", "null"]},
                    "files": {"type": "array", "items": _FILE_INFO_SCHEMA},
                    "skipped": {"type": "array", "items": _FILE_INFO_SCHEMA},
                    "error": {"type": "string"},
                },
                "required": ["files", "skipped"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_get_upcoming",
            description=(
                "Get upcoming calendar items and assignment due dates across your "
                "enrolled courses. Wraps Blackboard's calendar API — assignment "
                "due dates surface as items with type='GradebookColumn'. "
                "By default returns the next 2 weeks across every available course "
                "(server fans out per-course in parallel). Pass course_ids to scope "
                "to specific courses, since/until (ISO-8601) to override the window, "
                "or type to filter (e.g. type='GradebookColumn' for due dates only)."
            ),
            annotations={
                "title": "Get upcoming items / due dates",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "since": {
                        "type": "string",
                        "description": (
                            "ISO-8601 start of the window (e.g. '2026-05-09T00:00:00Z'). "
                            "Omit to default to now."
                        ),
                    },
                    "until": {
                        "type": "string",
                        "description": (
                            "ISO-8601 end of the window. Omit to default to "
                            "two weeks after `since`."
                        ),
                    },
                    "course_ids": {
                        "type": "array",
                        "description": (
                            "Optional list of course IDs to scope to. Omit to fan "
                            "out across all available enrolled courses."
                        ),
                        "items": {"type": "string", "pattern": _BB_ID_PATTERN},
                        "maxItems": _MAX_LIMIT,
                    },
                    "type": {
                        "type": "string",
                        "description": (
                            "Optional calendar item type filter. Use 'GradebookColumn' "
                            "for assignment due dates only."
                        ),
                        "enum": list(_CALENDAR_ITEM_TYPES),
                    },
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "items": {"type": "array", "items": _CALENDAR_ITEM_SCHEMA},
                    **_PAGINATION_OUTPUT_FIELDS,
                    "courseIdsQueried": {"type": "array", "items": {"type": "string"}},
                    "courseErrors": {
                        "type": "object",
                        "additionalProperties": {"type": "string"},
                    },
                },
                "required": ["items", "total", "count", "offset", "limit", "hasMore"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_get_announcements",
            description=(
                "Get announcements across your enrolled courses, newest first. "
                "By default fans out across every available course; pass "
                "course_ids=['_123_1'] to scope. Use since (ISO-8601) to filter "
                "to recent announcements only (e.g. \"this week\"). Each item "
                "includes the courseId it was posted to so cross-course views "
                "stay attributable. Supports pagination."
            ),
            annotations={
                "title": "Get announcements (cross-course)",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_ids": {
                        "type": "array",
                        "description": (
                            "Optional list of course IDs to scope to. Omit to fan "
                            "out across all available enrolled courses."
                        ),
                        "items": {"type": "string", "pattern": _BB_ID_PATTERN},
                        "maxItems": _MAX_LIMIT,
                    },
                    "since": {
                        "type": "string",
                        "description": (
                            "Optional ISO-8601 cutoff (e.g. '2026-05-09T00:00:00Z'). "
                            "Only announcements with `created` on/after this time "
                            "are returned. Filtered client-side after fetch."
                        ),
                    },
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "announcements": {"type": "array", "items": _ANNOUNCEMENT_SCHEMA},
                    **_PAGINATION_OUTPUT_FIELDS,
                    "courseIdsQueried": {"type": "array", "items": {"type": "string"}},
                    "courseErrors": {
                        "type": "object",
                        "additionalProperties": {"type": "string"},
                    },
                },
                "required": ["announcements", "total", "count", "offset", "limit", "hasMore"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_get_gradebook",
            description=(
                "Get gradebook columns across your enrolled courses, including your "
                "scores where available. By default fans out across every available "
                "course; pass course_ids=['_123_1'] to scope. Each column carries the "
                "courseId it belongs to so cross-course views stay attributable. "
                "Supports pagination."
            ),
            annotations={
                "title": "Get gradebook (cross-course)",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_ids": {
                        "type": "array",
                        "description": (
                            "Optional list of course IDs to scope to. Omit to fan "
                            "out across all available enrolled courses."
                        ),
                        "items": {"type": "string", "pattern": _BB_ID_PATTERN},
                        "maxItems": _MAX_LIMIT,
                    },
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "columns": {"type": "array", "items": _GRADEBOOK_COLUMN_SCHEMA},
                    **_PAGINATION_OUTPUT_FIELDS,
                    "gradesAvailable": {"type": "boolean"},
                    "gradeFetchError": {"type": ["string", "null"]},
                    "courseIdsQueried": {"type": "array", "items": {"type": "string"}},
                    "courseErrors": {
                        "type": "object",
                        "additionalProperties": {"type": "string"},
                    },
                },
                "required": [
                    "columns", "total", "count", "offset", "limit", "hasMore",
                    "gradesAvailable",
                ],
            },
        ),
        # ------------------------------------------------------------------
        # v0.3 tools — handlers live in ntulearn_mcp.handlers (lazy at call
        # time). Keep each entry's name/annotations/schema in lockstep with
        # the Handler contract in DEVELOPMENT-SPEC.md.
        # ------------------------------------------------------------------
        Tool(
            name=f"{_TOOL_PREFIX}_list_messages",
            description=(
                "List a user's Blackboard messages (mailbox) on NTULearn. "
                "Defaults to the inbox; pass folder='sent' for the outbox. "
                "Optionally filter to unread messages or a since (ISO-8601) "
                "cutoff. Supports pagination."
            ),
            annotations={
                "title": "List user messages",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "folder": _MESSAGE_FOLDER_SCHEMA,
                    "unread_only": {
                        "type": "boolean",
                        "description": "Return only unread messages when true.",
                        "default": False,
                    },
                    "since": _SINCE_SCHEMA,
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "messages": {"type": "array", "items": _MESSAGE_SCHEMA},
                    **_PAGINATION_OUTPUT_FIELDS,
                },
                "required": ["messages", "total", "count", "offset", "limit", "hasMore"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_read_message",
            description=(
                "Read a single Blackboard message by ID, including its subject, "
                "full body and recipient list. The ID comes from "
                "ntulearn_list_messages."
            ),
            annotations={
                "title": "Read a user message",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "message_id": {
                        "type": "string",
                        "description": "Blackboard message ID (e.g. _12345_1)",
                        "minLength": 1,
                        "maxLength": 200,
                        "pattern": _BB_ID_PATTERN,
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["message_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "message": _MESSAGE_SCHEMA,
                },
                "required": ["message"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_list_course_users",
            description=(
                "List the users enrolled in a course (instructors, teaching "
                "staff and students as the API exposes them). Supports "
                "pagination."
            ),
            annotations={
                "title": "List course users",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "users": {"type": "array", "items": _USER_SCHEMA},
                    **_PAGINATION_OUTPUT_FIELDS,
                },
                "required": ["users", "total", "count", "offset", "limit", "hasMore"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_list_course_groups",
            description=(
                "List the groups defined in a course (e.g. tutorial/lab "
                "groups). Supports pagination."
            ),
            annotations={
                "title": "List course groups",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "groups": {"type": "array", "items": _GROUP_SCHEMA},
                    **_PAGINATION_OUTPUT_FIELDS,
                },
                "required": ["groups", "total", "count", "offset", "limit", "hasMore"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_get_group_members",
            description=(
                "List the members of a specific course group (e.g. your "
                "tutorial group's roster)."
            ),
            annotations={
                "title": "Get course group members",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "group_id": {
                        "type": "string",
                        "description": "Group ID (e.g. _12345_1)",
                        "minLength": 1,
                        "maxLength": 200,
                        "pattern": _BB_ID_PATTERN,
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id", "group_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "users": {"type": "array", "items": _USER_SCHEMA},
                    "courseId": {"type": "string"},
                    "groupId": {"type": "string"},
                },
                "required": ["users"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_get_gradebook_attempts",
            description=(
                "List submission attempts for a gradebook column (assignment). "
                "Pass user_id to scope to one student's attempts. "
                "Supports pagination."
            ),
            annotations={
                "title": "Get gradebook attempts",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "column_id": {
                        "type": "string",
                        "description": "Gradebook column ID (e.g. _12345_1)",
                        "minLength": 1,
                        "maxLength": 200,
                        "pattern": _BB_ID_PATTERN,
                    },
                    "user_id": {
                        "type": "string",
                        "description": "Optional user ID to scope attempts to one student.",
                        "minLength": 1,
                        "maxLength": 200,
                        "pattern": _BB_ID_PATTERN,
                    },
                    "limit": _LIMIT_SCHEMA,
                    "offset": _OFFSET_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id", "column_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "attempts": {"type": "array", "items": _ATTEMPT_SCHEMA},
                    **_PAGINATION_OUTPUT_FIELDS,
                },
                "required": ["attempts", "total", "count", "offset", "limit", "hasMore"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_search_all_courses",
            description=(
                "Search across ALL enrolled courses' content trees for items "
                "matching a query. Like ntulearn_search_course_content but "
                "scoped to every available course at once, with per-item "
                "courseId and breadcrumb for attribution."
            ),
            annotations={
                "title": "Search content across all courses",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "query": _QUERY_SCHEMA,
                    "course_ids": _COURSE_IDS_SCHEMA,
                    "max_depth": {
                        **_MAX_DEPTH_SCHEMA,
                        "description": f"Maximum recursion depth per course (default 3, capped at {_MAX_DEPTH}).",
                        "default": 3,
                    },
                    "max_results": {
                        "type": "integer",
                        "description": f"Maximum number of matching items to return (default 50, capped at {_MAX_LIMIT})",
                        "default": 50,
                        "minimum": 1,
                        "maximum": _MAX_LIMIT,
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["query"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "matches": {
                        "type": "array",
                        "items": {
                            **_CONTENT_ITEM_SCHEMA,
                            "properties": {
                                **_CONTENT_ITEM_SCHEMA["properties"],
                                "courseId": {"type": ["string", "null"]},
                                "breadcrumb": {"type": "array", "items": {"type": "string"}},
                            },
                        },
                    },
                    "count": {"type": "integer"},
                    "courseIdsQueried": {"type": "array", "items": {"type": "string"}},
                    "courseErrors": {
                        "type": "object",
                        "additionalProperties": {"type": "string"},
                    },
                },
                "required": ["matches", "count"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_get_content_tree",
            description=(
                "Return a course's ENTIRE content tree as nested JSON "
                "(folders with children), unlike ntulearn_get_course_contents "
                "which returns one level per call. Use max_depth to bound the "
                "walk."
            ),
            annotations={
                "title": "Get full course content tree",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "max_depth": _MAX_DEPTH_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "courseId": {"type": "string"},
                    "tree": {
                        "type": "object",
                        "properties": {
                            "id": {"type": ["string", "null"]},
                            "title": {"type": ["string", "null"]},
                            "hasChildren": {"type": "boolean"},
                            "children": {"type": "array", "items": {"type": "object"}},
                        },
                    },
                    "count": {"type": "integer"},
                    "totalNodes": {"type": "integer"},
                },
                "required": ["courseId", "tree", "count", "totalNodes"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_download_course",
            description=(
                "Recursively download every file in a course's content tree "
                "to local disk, organised into course/subject folders. "
                "Skips already-downloaded files by default; pass "
                "skip_existing=false to re-download. Concurrency is bounded by "
                "`parallel`. Write tool (saves files on your machine)."
            ),
            annotations={
                "title": "Download an entire course to disk",
                "readOnlyHint": False,
                "destructiveHint": False,
                "idempotentHint": False,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "destination_dir": {
                        "type": "string",
                        "description": (
                            "Target root directory (absolute or ~-prefixed). "
                            "Defaults to ~/Downloads/NTU/<course>."
                        ),
                        "minLength": 1,
                        "maxLength": 1024,
                    },
                    "max_depth": _MAX_DEPTH_SCHEMA,
                    "include_extensions": {
                        "type": "string",
                        "description": (
                            "Optional comma-separated file extensions to download, "
                            "e.g. 'pdf,ppt,pptx,docx'. Omit to download all."
                        ),
                    },
                    "skip_existing": {
                        "type": "boolean",
                        "description": "Skip files already present at the destination.",
                        "default": True,
                    },
                    "parallel": {
                        "type": "integer",
                        "description": "Max concurrent downloads (1-16).",
                        "default": 4,
                        "minimum": 1,
                        "maximum": 16,
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "courseId": {"type": "string"},
                    "files": {
                        "type": "array",
                        "items": {
                            "type": "object",
                            "properties": {
                                "localPath": {"type": "string"},
                                "filename": {"type": "string"},
                                "sizeBytes": {"type": "integer"},
                                "courseFolder": {"type": "string"},
                            },
                        },
                    },
                    "skipped": {"type": "array", "items": {"type": "object"}},
                    "totalBytes": {"type": "integer"},
                    "destinationDir": {"type": "string"},
                },
                "required": ["courseId", "files", "totalBytes"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_whats_new",
            description=(
                "One-call 'what changed recently' digest: announcements, "
                "upcoming due dates and a gradebook summary across your "
                "courses since a cutoff. Defaults to the tracker's last-seen "
                "time (or the last 7 days). Set update_tracker=true to record "
                "`since` as the new last-seen for future calls."
            ),
            annotations={
                "title": "What's new across courses",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_ids": _COURSE_IDS_SCHEMA,
                    "since": {
                        "type": "string",
                        "description": (
                            "Optional ISO-8601 cutoff. Defaults to the tracker's "
                            "last-seen time, or 7 days ago if never set."
                        ),
                    },
                    "update_tracker": {
                        "type": "boolean",
                        "description": (
                            "Persist `since` as the new last-seen marker for "
                            "future ntulearn_whats_new calls."
                        ),
                        "default": False,
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "since": {"type": "string"},
                    "fetchedAt": {"type": "string"},
                    "announcements": {"type": "array", "items": _ANNOUNCEMENT_SCHEMA},
                    "upcoming": {"type": "array", "items": _CALENDAR_ITEM_SCHEMA},
                    "gradebookSummary": {
                        "type": ["object", "null"],
                        "properties": {
                            "columnCount": {"type": "integer"},
                            "columnsWithScore": {"type": "integer"},
                            "possibleTotal": {"type": ["number", "null"]},
                        },
                    },
                    "courseErrors": {
                        "type": "object",
                        "additionalProperties": {"type": "string"},
                    },
                },
                "required": ["since", "fetchedAt", "announcements", "upcoming"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_export_calendar_ics",
            description=(
                "Export calendar items (including assignment due dates) as an "
                "iCalendar (.ics) string you can paste into Google Calendar / "
                "Apple Calendar / Outlook. Optionally scope to course_ids and a "
                "since/until window."
            ),
            annotations={
                "title": "Export calendar as ICS",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_ids": _COURSE_IDS_SCHEMA,
                    "since": _SINCE_SCHEMA,
                    "until": _UNTIL_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "ics": {"type": "string"},
                    "itemCount": {"type": "integer"},
                    "supported": {"type": "boolean"},
                },
                "required": ["ics", "itemCount", "supported"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_export_gradebook_csv",
            description=(
                "Export your gradebook columns and scores across courses as a "
                "CSV string you can paste into a spreadsheet. Header: "
                "courseId,column,possible,score,grade,status."
            ),
            annotations={
                "title": "Export gradebook as CSV",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_ids": _COURSE_IDS_SCHEMA,
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "csv": {"type": "string"},
                    "columnCount": {"type": "integer"},
                    "courseCount": {"type": "integer"},
                },
                "required": ["csv", "columnCount", "courseCount"],
            },
        ),
        Tool(
            name=f"{_TOOL_PREFIX}_summarize_course",
            description=(
                "Bite-size briefing of a single course: title, description, "
                "instructors, enrollment count, upcoming due dates, recent "
                "announcements, a gradebook summary and the top-level content "
                "folders. Each sub-section degrades gracefully on error."
            ),
            annotations={
                "title": "Summarize a course",
                "readOnlyHint": True,
                "destructiveHint": False,
                "idempotentHint": True,
                "openWorldHint": True,
            },
            inputSchema={
                "type": "object",
                "properties": {
                    "course_id": _COURSE_ID_SCHEMA,
                    "include_contents": {
                        "type": "boolean",
                        "description": "Include the top-level content folders in the briefing.",
                        "default": True,
                    },
                    "response_format": _RESPONSE_FORMAT_SCHEMA,
                },
                "required": ["course_id"],
                "additionalProperties": False,
            },
            outputSchema={
                "type": "object",
                "properties": {
                    "courseId": {"type": "string"},
                    "title": {"type": ["string", "null"]},
                    "description": {"type": ["string", "null"]},
                    "instructors": {"type": "array", "items": {"type": "object"}},
                    "enrollmentCount": {"type": ["integer", "null"]},
                    "upcoming": {"type": "array", "items": _CALENDAR_ITEM_SCHEMA},
                    "recentAnnouncements": {"type": "array", "items": _ANNOUNCEMENT_SCHEMA},
                    "gradeSummary": {
                        "type": ["object", "null"],
                        "properties": {
                            "columnCount": {"type": "integer"},
                            "columnsWithScore": {"type": "integer"},
                            "possibleTotal": {"type": ["number", "null"]},
                        },
                    },
                    "contentTopFolders": {"type": "array", "items": _CONTENT_ITEM_SCHEMA},
                },
                "required": ["courseId"],
            },
        ),
    ]


# ---------------------------------------------------------------------------
# Tool handlers
# ---------------------------------------------------------------------------

@app.call_tool()
async def call_tool(
    name: str, arguments: dict[str, Any]
) -> tuple[list[TextContent | ImageContent], dict[str, Any]]:
    """Dispatch a tool call. Errors are raised so MCP wraps them with isError=True.

    On 401, swap in a fresh cookie once and retry. If retry still fails, the
    exception propagates and the SDK marks the result as an error.

    Returns a (unstructured_content, structured_content) tuple — both forms
    are propagated by the MCP framework so clients can use whichever they
    prefer (and the structured form is validated against any outputSchema
    declared on the tool).
    """
    try:
        return await _dispatch(name, arguments)
    except BbRouterExpiredError:
        await _refresh_client()
        return await _dispatch(name, arguments)


# Tool names implemented in ntulearn_mcp.handlers (WT-B). server.py has their
# Tool definitions but imports the handler module LAZILY at call time so this
# module still imports cleanly (and the existing server-local tools keep
# working) before handlers.py exists in a given build. Do not import handlers
# at module top.
_HANDLER_TOOL_NAMES: frozenset[str] = frozenset([
    f"{_TOOL_PREFIX}_list_messages",
    f"{_TOOL_PREFIX}_read_message",
    f"{_TOOL_PREFIX}_list_course_users",
    f"{_TOOL_PREFIX}_list_course_groups",
    f"{_TOOL_PREFIX}_get_group_members",
    f"{_TOOL_PREFIX}_get_gradebook_attempts",
    f"{_TOOL_PREFIX}_search_all_courses",
    f"{_TOOL_PREFIX}_get_content_tree",
    f"{_TOOL_PREFIX}_download_course",
    f"{_TOOL_PREFIX}_whats_new",
    f"{_TOOL_PREFIX}_export_calendar_ics",
    f"{_TOOL_PREFIX}_export_gradebook_csv",
    f"{_TOOL_PREFIX}_summarize_course",
])


def _load_handler(tool_name: str):
    """Lazily import ntulearn_mcp.handlers and return handle_<x> for a tool.

    Raises a clear RuntimeError when handlers.py is not present yet, so an MCP
    host that sees the tool listed gets an actionable error instead of a
    confusing ImportError/NameError.
    """
    try:
        from ntulearn_mcp import handlers
    except ImportError as e:
        raise RuntimeError(
            f"Tool {tool_name} is registered but its implementation is not "
            "available: ntulearn_mcp.handlers is missing (ImportError: "
            f"{e})."
        ) from e
    handler_name = "handle_" + tool_name[len(_TOOL_PREFIX) + 1:]
    handler = getattr(handlers, handler_name, None)
    if handler is None:
        raise RuntimeError(
            f"Tool {tool_name} is registered but handlers.{handler_name} "
            "was not found in ntulearn_mcp.handlers."
        )
    return handler


async def _dispatch(
    name: str, arguments: dict[str, Any]
) -> tuple[list[TextContent | ImageContent], dict[str, Any]]:
    client = get_client()
    handlers = {
        f"{_TOOL_PREFIX}_list_courses": _list_courses,
        f"{_TOOL_PREFIX}_get_course_contents": _get_course_contents,
        f"{_TOOL_PREFIX}_search_course_content": _search_course_content,
        f"{_TOOL_PREFIX}_download_file": _download_file,
        f"{_TOOL_PREFIX}_read_file_content": _read_file_content,
        f"{_TOOL_PREFIX}_get_upcoming": _get_upcoming,
        f"{_TOOL_PREFIX}_get_announcements": _get_announcements,
        f"{_TOOL_PREFIX}_get_gradebook": _get_gradebook,
    }
    handler = handlers.get(name)
    if handler is not None:
        return await handler(client, arguments)

    if name in _HANDLER_TOOL_NAMES:
        return await _load_handler(name)(client, arguments)

    known = sorted(set(handlers) | set(_HANDLER_TOOL_NAMES))
    raise ValueError(f"Unknown tool: {name}. Available: {known}")


# ---------------------------------------------------------------------------
# MCP resources (course briefings) + prompts (v0.3)
# ---------------------------------------------------------------------------

_RESOURCE_URI_PREFIX = "ntulearn://courses/"


@app.list_resource_templates()
async def list_resource_templates() -> list[ResourceTemplate]:
    """Expose the dynamic course-briefing URI template."""
    return [
        ResourceTemplate(
            name="ntulearn-course-briefing",
            title="NTULearn course briefing",
            uriTemplate="ntulearn://courses/{course_id}",
            description=(
                "JSON course briefing for a single NTULearn course (content "
                "rebuilt from the ntulearn_summarize_course tool). Replace "
                "{course_id} with a Blackboard course ID such as _12345_1."
            ),
            mimeType="application/json",
        )
    ]


@app.list_resources()
async def list_resources() -> list[Resource]:
    """Enumerate concrete course-briefing resources for enrolled courses.

    Returns an empty list when no cookie/client is available (e.g. during
    first-run inspection) — the URI is still resolvable through the template
    above, which is the primary discovery path for dynamic resources.
    """
    try:
        client = get_client()
    except Exception:
        return []
    try:
        enrollments = await client.get_my_enrollments()
    except Exception:
        return []
    resources = []
    for enrollment in enrollments:
        course_id = enrollment.get("courseId")
        if not course_id:
            continue
        resources.append(
            Resource(
                name=f"ntulearn-course-{course_id}",
                title=f"NTULearn course briefing: {course_id}",
                uri=f"{_RESOURCE_URI_PREFIX}{course_id}",
                description="JSON course briefing.",
                mimeType="application/json",
            )
        )
    return resources


@app.read_resource()
async def read_resource(uri: Any) -> ReadResourceResult:
    """Serve ntulearn://courses/{course_id} with a JSON course briefing.

    The briefing is produced by lazily importing ntulearn_mcp.handlers and
    calling handle_summarize_course — the same logic the tool runs — and
    embedding its payload as the resource's text content.
    """
    raw = str(uri)
    if not raw.startswith(_RESOURCE_URI_PREFIX):
        raise ValueError(
            f"Unsupported resource URI: {raw}. Expected "
            f"{_RESOURCE_URI_PREFIX}{{course_id}}"
        )
    course_id = raw[len(_RESOURCE_URI_PREFIX):].rstrip("/")
    if not course_id:
        raise ValueError(
            f"Missing course_id in resource URI: {raw}. Expected "
            f"{_RESOURCE_URI_PREFIX}{{course_id}}"
        )
    if not re.match(_BB_ID_PATTERN, course_id):
        raise ValueError(
            f"Invalid course_id {course_id!r} in resource URI {raw}."
        )

    try:
        from ntulearn_mcp import handlers
    except ImportError as e:
        raise RuntimeError(
            f"Cannot read resource {raw}: ntulearn_mcp.handlers is missing "
            f"(ImportError: {e})."
        ) from e

    _, payload = await handlers.handle_summarize_course(
        get_client(), {"course_id": course_id, "response_format": "json"}
    )
    # The handler always includes a (possibly empty) courseErrors key; only a
    # non-empty list is a failure worth surfacing.
    if payload.get("courseErrors"):
        raise ValueError(
            f"Course {course_id} could not be summarised: {payload['courseErrors']}"
        )
    return ReadResourceResult(
        contents=[
            TextResourceContents(
                uri=raw,
                mimeType="application/json",
                text=json.dumps(payload, indent=2),
            )
        ]
    )


# ---------------------------------------------------------------------------
# Prompt helpers
# ---------------------------------------------------------------------------

def _iso_cutoff(days_ago: int) -> str:
    """Return a UTC ISO-8601 timestamp `days_ago` from now, e.g. 2026-05-09T12:00:00Z."""
    stamp = datetime.now(timezone.utc) - timedelta(days=days_ago)
    return stamp.strftime("%Y-%m-%dT%H:%M:%SZ")


def _iso_future(days_ahead: int) -> str:
    """Return a UTC ISO-8601 timestamp `days_ahead` from now, e.g. 2026-05-16T12:00:00Z."""
    stamp = datetime.now(timezone.utc) + timedelta(days=days_ahead)
    return stamp.strftime("%Y-%m-%dT%H:%M:%SZ")


def _prompt_course_arg(courses: str | None) -> str:
    """Render the course_ids tool argument from a user-supplied courses string."""
    if not courses:
        return ""
    ids = [c.strip() for c in str(courses).split(",") if c.strip()]
    if not ids:
        return ""
    return f", course_ids={ids!r}"


@app.list_prompts()
async def list_prompts() -> list[Prompt]:
    """Advertise the two built-in prompt templates."""
    return [
        Prompt(
            name="ntulearn-weekly-brief",
            title="Weekly NTULearn brief",
            description=(
                "Generate a weekly digest of announcements and upcoming due "
                "dates across your courses. Chains ntulearn_get_announcements "
                "and ntulearn_get_upcoming with a computed since/until window."
            ),
            arguments=[
                PromptArgument(
                    name="courses",
                    description=(
                        "Optional comma-separated course IDs to scope to. "
                        "Omit for all enrolled courses."
                    ),
                    required=False,
                ),
                PromptArgument(
                    name="days",
                    description="Days back/forward in the window (default 7).",
                    required=False,
                ),
            ],
        ),
        Prompt(
            name="ntulearn-assignment-triage",
            title="NTULearn assignment triage",
            description=(
                "Prioritise your upcoming assignment due dates over the next "
                "N days. Chains ntulearn_get_upcoming (type='GradebookColumn') "
                "plus a grades lookup to suggest what to tackle first."
            ),
            arguments=[
                PromptArgument(
                    name="courses",
                    description=(
                        "Optional comma-separated course IDs to scope to. "
                        "Omit for all enrolled courses."
                    ),
                    required=False,
                ),
                PromptArgument(
                    name="days",
                    description="How far ahead to look for due dates (default 14).",
                    required=False,
                ),
            ],
        ),
    ]


@app.get_prompt()
async def get_prompt(
    name: str, arguments: dict[str, Any] | None = None
) -> GetPromptResult:
    """Return the prompt TEXT for a named template.

    Prompts only instruct the model which tools to chain and with what
    arguments — no tools are executed here. Recommended calls are embedded
    so the model has exact since/until values.
    """
    args = arguments or {}
    courses = args.get("courses")
    if name == "ntulearn-weekly-brief":
        try:
            days = int(args.get("days", 7))
        except (TypeError, ValueError):
            raise ValueError("days must be an integer") from None
        if days < 1:
            raise ValueError("days must be >= 1")
        since = _iso_cutoff(days)
        until = _iso_future(0)
        course_arg = _prompt_course_arg(courses)
        text = (
            "Produce a weekly NTULearn brief covering the last/looking "
            f"{days} days. Run these tools and combine the results:" + "\n"
            f"1. ntulearn_get_announcements(since='{since}'{course_arg}) "
            "— recent announcements." + "\n"
            f"2. ntulearn_get_upcoming(since='{since}', until='{until}'{course_arg})"
            " — upcoming calendar items and due dates." + "\n"
            "Summarise: what's new in each course, what's due, and any "
            "deadlines the user should watch in the next few days."
        )
        return GetPromptResult(
            description="Weekly NTULearn brief",
            messages=[PromptMessage(role="user", content=TextContent(type="text", text=text))],
        )
    if name == "ntulearn-assignment-triage":
        try:
            days = int(args.get("days", 14))
        except (TypeError, ValueError):
            raise ValueError("days must be an integer") from None
        if days < 1:
            raise ValueError("days must be >= 1")
        until = _iso_future(days)
        course_arg = _prompt_course_arg(courses)
        gradebook_arg = ""
        ids = [c.strip() for c in str(args.get("courses", "")).split(",") if c.strip()]
        if ids:
            gradebook_arg = f"course_ids={ids!r}"
        text = (
            "Running an assignment triage for the next "
            f"{days} days. Run these tools and combine the results:" + "\n"
            f"1. ntulearn_get_upcoming(type='GradebookColumn', until='{until}'{course_arg})"
            " — assignment due dates." + "\n"
            "2. ntulearn_get_gradebook(" + gradebook_arg + ") — current scores "
            "for context." + "\n"
            "For each due assignment, list: course, title, due date, how much "
            "time is left, and a recommended order of attack (closest deadline "
            "first, biggest weight first)."
        )
        return GetPromptResult(
            description="NTULearn assignment triage",
            messages=[PromptMessage(role="user", content=TextContent(type="text", text=text))],
        )
    raise ValueError(f"Unknown prompt: {name}")

# ---------------------------------------------------------------------------
# Markdown rendering helpers
# ---------------------------------------------------------------------------

def _md_courses(courses: list[dict[str, Any]], meta: dict[str, Any]) -> str:
    lines = [f"# Courses ({meta['total']} total)", ""]
    if not courses:
        lines.append("_No courses to show._")
    else:
        for c in courses:
            last = c.get("lastAccessed") or "—"
            lines.append(
                f"- **{c.get('title', '?')}** "
                f"`{c.get('courseId', '?')}` "
                f"· available={c.get('available', '?')} "
                f"· last accessed {last}"
            )
    lines.append(_md_pagination_footer(meta))
    return "\n".join(lines)


def _md_content_items(items: list[dict[str, Any]], meta: dict[str, Any]) -> str:
    lines = [f"# Content items ({meta['total']} total)", ""]
    if not items:
        lines.append("_No items._")
    else:
        for it in items:
            arrow = "📁" if it.get("hasChildren") else "📄"
            lines.append(
                f"- {arrow} **{it.get('title', '?')}** "
                f"`{it.get('id', '?')}` "
                f"· handler={it.get('contentHandlerId', '?')}"
            )
    lines.append(_md_pagination_footer(meta))
    return "\n".join(lines)


def _md_upcoming(
    items: list[dict[str, Any]],
    meta: dict[str, Any],
    course_errors: dict[str, str],
) -> str:
    lines = [f"# Upcoming ({meta['total']} total)", ""]
    if course_errors:
        lines.append(
            f"_Note: {len(course_errors)} course(s) returned errors and were skipped._"
        )
        lines.append("")
    if not items:
        lines.append("_Nothing scheduled in the window._")
    else:
        lines.append("| When | Title | Type | Course | Gradable |")
        lines.append("|---|---|---|---|---|")
        for it in items:
            start = it.get("start") or "—"
            end = it.get("end")
            when = f"{start} → {end}" if end and end != start else start
            gradable = it.get("gradable")
            gradable_str = "Yes" if gradable else ("No" if gradable is False else "—")
            lines.append(
                f"| {when} "
                f"| {it.get('title') or '?'} "
                f"| {it.get('type') or '—'} "
                f"| {it.get('courseId') or '—'} "
                f"| {gradable_str} |"
            )
    lines.append(_md_pagination_footer(meta))
    return "\n".join(lines)


def _md_announcements(items: list[dict[str, Any]], meta: dict[str, Any]) -> str:
    lines = [f"# Announcements ({meta['total']} total)", ""]
    if not items:
        lines.append("_No announcements._")
    else:
        for a in items:
            created = a.get("created") or "—"
            course_id = a.get("courseId")
            header = f"## {a.get('title', '?')}  ·  {created}"
            if course_id:
                header += f"  ·  `{course_id}`"
            lines.append(header)
            body = (a.get("body") or "").strip()
            lines.append(body if body else "_(no body)_")
            lines.append("")
    lines.append(_md_pagination_footer(meta))
    return "\n".join(lines)


def _md_gradebook(columns: list[dict[str, Any]], meta: dict[str, Any], grades_available: bool, error: str | None) -> str:
    lines = [f"# Gradebook ({meta['total']} columns)", ""]
    if not grades_available:
        lines.append(f"_Grades not available: {error}_")
        lines.append("")
    if not columns:
        lines.append("_No columns._")
    else:
        lines.append("| Course | Column | Possible | Score | Grade | Status |")
        lines.append("|---|---|---|---|---|---|")
        for c in columns:
            lines.append(
                f"| `{c.get('courseId', '—')}` "
                f"| {c.get('displayName') or c.get('name') or '?'} "
                f"| {c.get('possible', '—')} "
                f"| {c.get('score', '—')} "
                f"| {c.get('grade', '—')} "
                f"| {c.get('status', '—')} |"
            )
    lines.append(_md_pagination_footer(meta))
    return "\n".join(lines)


def _md_search_results(matches: list[dict[str, Any]]) -> str:
    lines = [f"# Search matches ({len(matches)})", ""]
    if not matches:
        lines.append("_No matches._")
    else:
        for m in matches:
            crumb = " › ".join(m.get("breadcrumb") or [])
            lines.append(f"- **{m.get('title', '?')}** `{m.get('id', '?')}`  ")
            lines.append(f"  _{crumb}_")
    return "\n".join(lines)


def _md_files(payload: dict[str, Any], heading: str) -> str:
    lines = [f"# {heading}", "", f"**Item:** {payload.get('title', '?')}  ", ""]
    files = payload.get("files") or []
    skipped = payload.get("skipped") or []
    if files:
        lines.append("## Files")
        for f in files:
            if "localPath" in f:
                lines.append(
                    f"- `{f['filename']}` ({_format_bytes(f.get('sizeBytes', 0))}) "
                    f"→ `{f['localPath']}`"
                )
            elif "text" in f:
                count_bits = []
                if f.get("pageCount"):
                    count_bits.append(f"{f['pageCount']} pages")
                if f.get("slideCount"):
                    count_bits.append(f"{f['slideCount']} slides")
                if f.get("sheetCount"):
                    count_bits.append(f"{f['sheetCount']} sheets")
                counts = f" · {', '.join(count_bits)}" if count_bits else ""
                lines.append(
                    f"### {f['filename']} ({f['kind']}{counts}, "
                    f"{_format_bytes(f.get('sizeBytes', 0))})"
                )
                text = f.get("text", "").strip()
                if not text:
                    lines.append("_(empty)_")
                elif len(text) > 5000:
                    lines.append(text[:5000] + "\n…_(truncated in markdown view; use response_format='json' for full text)_")
                else:
                    lines.append(text)
            else:
                lines.append(f"- `{f.get('filename', '?')}` (url: {f.get('url', '?')})")
        lines.append("")
    if skipped:
        lines.append("## Skipped")
        for s in skipped:
            lines.append(f"- `{s['filename']}`: {s['reason']}")
    if not files and not skipped:
        lines.append("_(nothing to show)_")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# Individual tool implementations
# ---------------------------------------------------------------------------

async def _list_courses(
    client: NTULearnClient, args: dict[str, Any]
) -> tuple[list[TextContent], dict[str, Any]]:
    include_disabled = bool(args.get("include_disabled", False))
    offset, limit = _resolve_pagination_args(args)
    fmt = _resolve_response_format(args)

    enrollments = await client.get_my_enrollments()
    if not include_disabled:
        enrollments = [
            e for e in enrollments
            if (e.get("availability") or {}).get("available") == "Yes"
        ]

    course_ids = [e["courseId"] for e in enrollments]
    if not course_ids:
        _, meta = _slice_with_pagination([], offset, limit)
        payload = {"courses": [], **meta}
        text = _md_courses([], meta) if fmt == "markdown" else None
        return _emit(payload, text)

    last_accessed_map = {e["courseId"]: e.get("lastAccessed") for e in enrollments}
    availability_map = {
        e["courseId"]: (e.get("availability") or {}).get("available", "Unknown")
        for e in enrollments
    }

    courses_raw = await client.get_courses_batch(course_ids)
    rows = []
    for course in courses_raw:
        cid = course.get("id")
        rows.append({
            "courseId": cid,
            "title": course.get("name") or course.get("displayName") or course.get("id"),
            "available": availability_map.get(cid, "Unknown"),
            "lastAccessed": last_accessed_map.get(cid),
        })
    rows.sort(key=lambda c: c["lastAccessed"] or "", reverse=True)

    page, meta = _slice_with_pagination(rows, offset, limit)
    payload = {"courses": page, **meta}
    text = _md_courses(page, meta) if fmt == "markdown" else None
    return _emit(payload, text)


async def _get_course_contents(
    client: NTULearnClient, args: dict[str, Any]
) -> tuple[list[TextContent], dict[str, Any]]:
    course_id = args["course_id"]
    parent_id = args.get("parent_id")
    offset, limit = _resolve_pagination_args(args)
    fmt = _resolve_response_format(args)

    if parent_id:
        items = await client.get_content_children(course_id, str(parent_id))
    else:
        items = await client.get_course_contents(course_id)
    stripped = [_strip_content(item) for item in items]
    page, meta = _slice_with_pagination(stripped, offset, limit)
    payload = {"items": page, **meta}
    text = _md_content_items(page, meta) if fmt == "markdown" else None
    return _emit(payload, text)


async def _search_course_content(
    client: NTULearnClient, args: dict[str, Any]
) -> tuple[list[TextContent], dict[str, Any]]:
    course_id = args["course_id"]
    query = str(args["query"]).strip().lower()
    if not query:
        raise ValueError("search query cannot be blank.")

    max_depth = int(args.get("max_depth", 5))
    max_results = int(args.get("max_results", 50))
    if max_depth < 1 or max_depth > _MAX_DEPTH:
        raise ValueError(f"max_depth must be 1..{_MAX_DEPTH}")
    if max_results < 1 or max_results > _MAX_LIMIT:
        raise ValueError(f"max_results must be 1..{_MAX_LIMIT}")
    fmt = _resolve_response_format(args)

    matches: list[dict[str, Any]] = []
    visited: set[str] = set()
    semaphore = asyncio.Semaphore(5)

    async def walk(items: list[dict[str, Any]], path: list[str], depth: int) -> None:
        if depth > max_depth or len(matches) >= max_results:
            return

        child_tasks = []
        for item in items:
            if len(matches) >= max_results:
                break

            item_id = item.get("id")
            if item_id:
                if item_id in visited:
                    continue
                visited.add(item_id)

            title = item.get("title") or ""
            desc_raw = item.get("description") or {}
            desc = (desc_raw.get("rawText") if isinstance(desc_raw, dict) else desc_raw) or ""

            current_path = path + [title]

            if query in title.lower() or query in desc.lower():
                stripped = _strip_content(item)
                stripped["breadcrumb"] = current_path
                matches.append(stripped)

            if item.get("hasChildren") and len(matches) < max_results:
                async def fetch_children(i=item, p=current_path):
                    async with semaphore:
                        children = await client.get_content_children(course_id, i["id"])
                    await walk(children, p, depth + 1)
                child_tasks.append(fetch_children())

        if child_tasks:
            await asyncio.gather(*child_tasks)

    top_level = await client.get_course_contents(course_id)
    await walk(top_level, [], 0)

    payload = {"matches": matches, "count": len(matches)}
    text = _md_search_results(matches) if fmt == "markdown" else None
    return _emit(payload, text)


async def _download_file(
    client: NTULearnClient, args: dict[str, Any]
) -> tuple[list[TextContent], dict[str, Any]]:
    course_id = args["course_id"]
    content_id = args["content_id"]
    fmt = _resolve_response_format(args)
    dest_dir = _resolve_destination_dir(args.get("destination_dir"))

    item, handler_id, pairs = await _resolve_content_files(client, course_id, content_id)

    if not pairs:
        payload = {
            "contentId": content_id,
            "title": item.get("title"),
            "contentHandlerId": handler_id,
            "files": [],
            "destinationDir": str(dest_dir),
            "error": "No download URL found. Content handler type may not be supported.",
        }
        text = (
            f"# Nothing downloaded\n\nItem **{item.get('title', '?')}** "
            f"(handler `{handler_id}`) has no resolvable file links."
        ) if fmt == "markdown" else None
        return _emit(payload, text)

    dest_dir.mkdir(parents=True, exist_ok=True)

    def _sanitize(name: str) -> str:
        return re.sub(r'[\\/*?:"<>|]', "_", name)

    def _deduplicate(name: str) -> str:
        candidate = name
        stem, dot, ext = name.rpartition(".")
        base = stem if dot else name
        suffix = ext if dot else ""
        n = 2
        while candidate in used_names or (dest_dir / candidate).exists():
            candidate = f"{base} ({n}).{suffix}" if suffix else f"{base} ({n})"
            n += 1
        return candidate

    used_names: set[str] = set()
    saved: list[dict[str, Any]] = []

    for url, detected_filename in pairs:
        filename = detected_filename
        if not filename:
            url_path = url.split("?")[0]
            filename = unquote(url_path.split("/")[-1]) or "download"
        filename = _sanitize(filename)
        filename = _deduplicate(filename)
        used_names.add(filename)

        dest = dest_dir / filename
        try:
            content_bytes, _ = await client.download_bytes(url)
        except BbRouterExpiredError:
            client = await _refresh_client()
            content_bytes, _ = await client.download_bytes(url)
        dest.write_bytes(content_bytes)

        saved.append({
            "localPath": str(dest.resolve()),
            "filename": filename,
            "sizeBytes": len(content_bytes),
        })

    payload = {
        "contentId": content_id,
        "title": item.get("title"),
        "files": saved,
        "destinationDir": str(dest_dir.resolve()),
    }
    text = _md_files(payload, "Files downloaded") if fmt == "markdown" else None
    return _emit(payload, text)


def _resolve_destination_dir(raw: Any) -> Path:
    """Resolve the download destination directory.

    Precedence: explicit `destination_dir` arg → `NTULEARN_DOWNLOAD_DIR` env →
    `DOWNLOAD_DIR` module default (`./downloads/`).

    Accepts absolute paths or `~`-prefixed paths. Relative paths are resolved
    against `DOWNLOAD_DIR.parent` (the server's CWD by default) so the model
    cannot accidentally write outside the user's intended root by passing
    `subdir/foo`.
    """
    if raw is not None:
        if not isinstance(raw, str):
            raise ValueError("destination_dir must be a string")
        candidate = raw.strip()
        if not candidate:
            raise ValueError("destination_dir cannot be empty")
        return Path(candidate).expanduser()
    env_val = os.environ.get("NTULEARN_DOWNLOAD_DIR")
    if env_val:
        return Path(env_val).expanduser()
    return DOWNLOAD_DIR


async def _read_file_content(
    client: NTULearnClient, args: dict[str, Any]
) -> tuple[list[TextContent | ImageContent], dict[str, Any]]:
    """Resolve files attached to a content item, fetch bytes, return content inline.

    Bypasses the local-filesystem hop that breaks Claude Desktop's sandbox:
    rather than writing to ./downloads, the bytes are extracted in-process and
    returned as MCP content blocks. PDFs default to vision mode (text + page
    images via PyMuPDF, matching Claude.ai's native PDF flow); pass mode='text'
    for the cheaper pypdf-only path. Office formats and text-likes are always
    text-only.
    """
    course_id = args["course_id"]
    content_id = args["content_id"]
    fmt = _resolve_response_format(args)
    pdf_mode = _resolve_pdf_mode(args)
    pages_filter = _parse_page_range(args.get("pages"))

    item, handler_id, pairs = await _resolve_content_files(client, course_id, content_id)

    if not pairs:
        payload = {
            "contentId": content_id,
            "title": item.get("title"),
            "contentHandlerId": handler_id,
            "files": [],
            "skipped": [],
            "error": "No download URL found. Content handler type may not be supported.",
        }
        text = (
            f"# No content\n\nItem **{item.get('title', '?')}** "
            f"(handler `{handler_id}`) has no resolvable file links."
        ) if fmt == "markdown" else None
        return _emit(payload, text)

    files_out: list[dict[str, Any]] = []
    skipped: list[dict[str, Any]] = []
    image_blocks: list[ImageContent] = []
    total_bytes = 0

    for url, detected_filename in pairs:
        filename = detected_filename
        if not filename:
            url_path = url.split("?")[0]
            filename = unquote(url_path.split("/")[-1]) or "download"

        if total_bytes >= _MAX_TOTAL_BYTES:
            skipped.append({
                "filename": filename,
                "reason": (
                    f"Skipped: cumulative batch size already exceeds "
                    f"{_format_bytes(_MAX_TOTAL_BYTES)} cap. Use ntulearn_download_file."
                ),
            })
            continue

        try:
            content_bytes, content_type = await client.download_bytes(url)
        except BbRouterExpiredError:
            client = await _refresh_client()
            content_bytes, content_type = await client.download_bytes(url)

        size = len(content_bytes)
        if size > _MAX_FILE_BYTES:
            skipped.append({
                "filename": filename,
                "reason": (
                    f"File too large ({_format_bytes(size)} > "
                    f"{_format_bytes(_MAX_FILE_BYTES)} cap). Use ntulearn_download_file."
                ),
                "sizeBytes": size,
                "contentType": content_type,
            })
            continue

        if total_bytes + size > _MAX_TOTAL_BYTES:
            skipped.append({
                "filename": filename,
                "reason": (
                    f"Skipped: would exceed batch cap of "
                    f"{_format_bytes(_MAX_TOTAL_BYTES)}. Use ntulearn_download_file."
                ),
                "sizeBytes": size,
                "contentType": content_type,
            })
            continue

        total_bytes += size
        if _classify_kind(filename, content_type) == "pdf" and pdf_mode == "vision":
            entry = _extract_pdf_vision(
                filename, content_type, content_bytes, size, pages_filter
            )
        else:
            entry = _extract_content(filename, content_type, content_bytes)
        if entry.get("kind") == "binary":
            skipped.append({
                "filename": filename,
                "reason": entry["error"],
                "sizeBytes": size,
                "contentType": content_type,
            })
            continue

        # Pull rendered images (if any) out of the structured payload and
        # convert to ImageContent blocks. The bytes are too large to round-trip
        # through structured content, and clients want them as their own blocks
        # anyway so the model receives them as vision input.
        for label, png_bytes in entry.pop("_images", []):
            image_blocks.append(ImageContent(
                type="image",
                data=base64.b64encode(png_bytes).decode("ascii"),
                mimeType="image/png",
                annotations=None,
            ))
        files_out.append(entry)

    payload = {
        "contentId": content_id,
        "title": item.get("title"),
        "files": files_out,
        "skipped": skipped,
    }
    text = _md_files(payload, "File contents") if fmt == "markdown" else None
    text_blocks, structured = _emit(payload, text)
    return [*text_blocks, *image_blocks], structured


async def _resolve_enrolled_course_ids(
    client: NTULearnClient, *, include_disabled: bool = False
) -> list[str]:
    """Return the list of course IDs the current user is enrolled in.

    Used by cross-course aggregators (get_upcoming, get_announcements,
    get_gradebook) when no explicit course_ids list is supplied. By default
    excludes courses where availability is not 'Yes' so the fan-out doesn't
    spend requests on disabled / archived terms.
    """
    enrollments = await client.get_my_enrollments()
    if not include_disabled:
        enrollments = [
            e for e in enrollments
            if (e.get("availability") or {}).get("available") == "Yes"
        ]
    return [e["courseId"] for e in enrollments if e.get("courseId")]


async def _get_upcoming(
    client: NTULearnClient, args: dict[str, Any]
) -> tuple[list[TextContent], dict[str, Any]]:
    raw_since = args.get("since")
    raw_until = args.get("until")
    since = _validate_iso8601(raw_since, name="since") if raw_since is not None else None
    until = _validate_iso8601(raw_until, name="until") if raw_until is not None else None

    item_type = args.get("type")
    if item_type is not None and item_type not in _CALENDAR_ITEM_TYPES:
        raise ValueError(
            f"type must be one of {list(_CALENDAR_ITEM_TYPES)}; got {item_type!r}"
        )

    raw_course_ids = args.get("course_ids")
    if raw_course_ids is None or (
        isinstance(raw_course_ids, list) and not raw_course_ids
    ):
        course_ids = await _resolve_enrolled_course_ids(client)
    else:
        if not isinstance(raw_course_ids, list):
            raise ValueError("course_ids must be a list of strings")
        course_ids = [str(cid) for cid in raw_course_ids]

    offset, limit = _resolve_pagination_args(args)
    fmt = _resolve_response_format(args)

    if not course_ids:
        _, meta = _slice_with_pagination([], offset, limit)
        payload = {
            "items": [],
            **meta,
            "courseIdsQueried": [],
            "courseErrors": {},
        }
        text = _md_upcoming([], meta, {}) if fmt == "markdown" else None
        return _emit(payload, text)

    tasks = [
        client.get_calendar_items(
            course_id=cid, since=since, until=until, item_type=item_type
        )
        for cid in course_ids
    ]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    items: list[dict[str, Any]] = []
    course_errors: dict[str, str] = {}
    for cid, result in zip(course_ids, results):
        if isinstance(result, BbRouterExpiredError):
            raise result
        if isinstance(result, Exception):
            course_errors[cid] = str(result)
            continue
        for raw in result:
            items.append(_strip_calendar_item(raw, cid))

    items.sort(key=lambda x: x.get("start") or "￿")

    page, meta = _slice_with_pagination(items, offset, limit)
    payload: dict[str, Any] = {
        "items": page,
        **meta,
        "courseIdsQueried": course_ids,
        "courseErrors": course_errors,
    }
    text = _md_upcoming(page, meta, course_errors) if fmt == "markdown" else None
    return _emit(payload, text)


async def _get_announcements(
    client: NTULearnClient, args: dict[str, Any]
) -> tuple[list[TextContent], dict[str, Any]]:
    raw_since = args.get("since")
    since = _validate_iso8601(raw_since, name="since") if raw_since is not None else None

    raw_course_ids = args.get("course_ids")
    if raw_course_ids is None or (
        isinstance(raw_course_ids, list) and not raw_course_ids
    ):
        course_ids = await _resolve_enrolled_course_ids(client)
    else:
        if not isinstance(raw_course_ids, list):
            raise ValueError("course_ids must be a list of strings")
        course_ids = [str(cid) for cid in raw_course_ids]

    offset, limit = _resolve_pagination_args(args)
    fmt = _resolve_response_format(args)

    if not course_ids:
        _, meta = _slice_with_pagination([], offset, limit)
        payload = {
            "announcements": [],
            **meta,
            "courseIdsQueried": [],
            "courseErrors": {},
        }
        text = _md_announcements([], meta) if fmt == "markdown" else None
        return _emit(payload, text)

    tasks = [client.get_announcements(cid) for cid in course_ids]
    results = await asyncio.gather(*tasks, return_exceptions=True)

    rows: list[dict[str, Any]] = []
    course_errors: dict[str, str] = {}
    for cid, result in zip(course_ids, results):
        if isinstance(result, BbRouterExpiredError):
            raise result
        if isinstance(result, Exception):
            course_errors[cid] = str(result)
            continue
        for a in result:
            body_raw = a.get("body") or {}
            # Blackboard stores announcement bodies as HTML in `body.rawText`. Strip
            # tags so callers see plain text, matching how the HTML file path in
            # _extract_content already behaves.
            body_html = body_raw.get("rawText") if isinstance(body_raw, dict) else body_raw
            created = a.get("created")
            if since is not None and created is not None and created < since:
                continue
            rows.append({
                "id": a.get("id"),
                "courseId": cid,
                "title": a.get("title"),
                "body": _strip_html(body_html),
                "created": created,
                "modified": a.get("modified"),
                "available": (a.get("availability") or {}).get("available"),
            })

    rows.sort(key=lambda r: r.get("created") or "", reverse=True)

    page, meta = _slice_with_pagination(rows, offset, limit)
    payload = {
        "announcements": page,
        **meta,
        "courseIdsQueried": course_ids,
        "courseErrors": course_errors,
    }
    text = _md_announcements(page, meta) if fmt == "markdown" else None
    return _emit(payload, text)


async def _get_gradebook(
    client: NTULearnClient, args: dict[str, Any]
) -> tuple[list[TextContent], dict[str, Any]]:
    raw_course_ids = args.get("course_ids")
    if raw_course_ids is None or (
        isinstance(raw_course_ids, list) and not raw_course_ids
    ):
        course_ids = await _resolve_enrolled_course_ids(client)
    else:
        if not isinstance(raw_course_ids, list):
            raise ValueError("course_ids must be a list of strings")
        course_ids = [str(cid) for cid in raw_course_ids]

    offset, limit = _resolve_pagination_args(args)
    fmt = _resolve_response_format(args)

    grades_available = True
    grade_fetch_error: str | None = None
    user_id: str | None = None
    try:
        user_id = await client.get_my_user_id()
    except BbRouterExpiredError:
        raise
    except Exception as e:
        grades_available = False
        grade_fetch_error = str(e)

    if not course_ids:
        _, meta = _slice_with_pagination([], offset, limit)
        payload = {
            "columns": [],
            **meta,
            "gradesAvailable": grades_available,
            "gradeFetchError": grade_fetch_error,
            "courseIdsQueried": [],
            "courseErrors": {},
        }
        text = (
            _md_gradebook([], meta, grades_available, grade_fetch_error)
            if fmt == "markdown"
            else None
        )
        return _emit(payload, text)

    async def fetch_one(
        cid: str,
    ) -> tuple[list[dict[str, Any]], list[dict[str, Any]], str | None]:
        """Per-course fetch. Returns (columns, grades, grade_error).

        Column fetch failures propagate (caught in the gather aggregation).
        Grade fetch failures are reported per-course rather than failing the
        whole call — the columns themselves are still useful, and only some
        courses may have grading enabled.
        """
        columns = await client.get_gradebook_columns(cid)
        if user_id is None:
            return columns, [], grade_fetch_error
        try:
            grades = await client.get_user_grades(cid, user_id)
            return columns, grades, None
        except BbRouterExpiredError:
            raise
        except Exception as e:
            return columns, [], str(e)

    results = await asyncio.gather(
        *(fetch_one(cid) for cid in course_ids), return_exceptions=True
    )

    columns_result: list[dict[str, Any]] = []
    course_errors: dict[str, str] = {}
    grade_errors: list[str] = []
    if grade_fetch_error:
        grade_errors.append(grade_fetch_error)
    for cid, result in zip(course_ids, results):
        if isinstance(result, BbRouterExpiredError):
            raise result
        if isinstance(result, Exception):
            course_errors[cid] = str(result)
            continue
        columns, grades_raw, per_course_grade_error = result
        if per_course_grade_error:
            grade_errors.append(per_course_grade_error)
        grade_map: dict[str, dict[str, Any]] = {
            g["columnId"]: g for g in grades_raw if "columnId" in g
        }
        for col in columns:
            col_id = col.get("id")
            score = col.get("score") or {}
            grade_entry = grade_map.get(col_id, {})
            columns_result.append({
                "id": col_id,
                "courseId": cid,
                "name": col.get("name"),
                "displayName": col.get("displayName"),
                "possible": score.get("possible"),
                "available": (col.get("availability") or {}).get("available"),
                "contentId": col.get("contentId"),
                "score": grade_entry.get("score"),
                "grade": grade_entry.get("grade"),
                "status": grade_entry.get("status"),
            })

    if grade_errors:
        # Any course-level grade fetch failure flips gradesAvailable to False
        # to match the v0.1.x single-course contract: callers checked this
        # flag to know whether `score`/`grade` fields are reliable.
        grades_available = False
        if grade_fetch_error is None:
            grade_fetch_error = "; ".join(dict.fromkeys(grade_errors))

    page, meta = _slice_with_pagination(columns_result, offset, limit)
    payload = {
        "columns": page,
        **meta,
        "gradesAvailable": grades_available,
        "gradeFetchError": grade_fetch_error,
        "courseIdsQueried": course_ids,
        "courseErrors": course_errors,
    }
    text = (
        _md_gradebook(page, meta, grades_available, grade_fetch_error)
        if fmt == "markdown"
        else None
    )
    return _emit(payload, text)


# ---------------------------------------------------------------------------
# Entrypoint
# ---------------------------------------------------------------------------

async def _run() -> None:
    async with stdio_server() as (read_stream, write_stream):
        await app.run(
            read_stream,
            write_stream,
            app.create_initialization_options(),
        )


def main() -> None:
    asyncio.run(_run())


if __name__ == "__main__":
    main()
